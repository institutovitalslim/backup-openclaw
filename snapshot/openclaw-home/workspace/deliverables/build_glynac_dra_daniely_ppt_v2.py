from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
from pathlib import Path

OUT = Path('/root/.openclaw/workspace/deliverables')
TMP = OUT / 'tmp_glynac_v2'
TMP.mkdir(parents=True, exist_ok=True)

photos = [
    Path('/root/.openclaw/media/inbound/file_186---e75a5ed2-8053-4cd0-bcdd-2b5ef5a70bc1.jpg'),
    Path('/root/.openclaw/media/inbound/file_187---3825409d-f26b-49ca-b763-4884eb4c4091.jpg'),
    Path('/root/.openclaw/media/inbound/file_188---379cfcdb-eb2d-4c8b-a2c8-487e1823d226.jpg'),
    Path('/root/.openclaw/media/inbound/file_189---7c6d471e-d5e7-4d9d-aa3e-a58397dc6a84.jpg'),
    Path('/root/.openclaw/media/inbound/file_190---221a1383-edf3-427e-bcef-1e6be5de8498.jpg'),
    Path('/root/.openclaw/media/inbound/file_191---76af0839-730b-44e9-bacd-b20082c8c29e.jpg'),
]

CREAM = RGBColor(247, 243, 236)
GREEN = RGBColor(13, 85, 71)
DEEP = RGBColor(17, 24, 22)
MOSS = RGBColor(186, 208, 199)
TEXT = RGBColor(23, 27, 25)
MUTED = RGBColor(94, 103, 99)
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(195, 163, 98)
LINE = RGBColor(217, 223, 219)

prs = Presentation()
prs.slide_width = Inches(8.333)
prs.slide_height = Inches(10.417)
SW, SH = prs.slide_width, prs.slide_height


def crop_fill(src, dst, size, dark=False):
    img = Image.open(src).convert('RGB')
    tw, th = size
    src_ratio = img.width / img.height
    dst_ratio = tw / th
    if src_ratio > dst_ratio:
        nw = int(img.height * dst_ratio)
        left = (img.width - nw) // 2
        img = img.crop((left, 0, left + nw, img.height))
    else:
        nh = int(img.width / dst_ratio)
        top = (img.height - nh) // 2
        img = img.crop((0, top, img.width, top + nh))
    img = img.resize((tw, th), Image.Resampling.LANCZOS)
    if dark:
        overlay = Image.new('RGBA', (tw, th), (0, 0, 0, 0))
        draw = ImageDraw.Draw(overlay)
        for y in range(th):
            alpha = int(165 * (y / th))
            draw.line((0, y, tw, y), fill=(7, 10, 9, alpha))
        for x in range(tw):
            alpha = int(130 * ((tw - x) / tw))
            draw.line((x, 0, x, th), fill=(7, 10, 9, alpha))
        img = Image.alpha_composite(img.convert('RGBA'), overlay).convert('RGB')
    img.save(dst, quality=95)


proc = []
for i, p in enumerate(photos, start=1):
    out = TMP / f'img_{i}.jpg'
    crop_fill(p, out, (1200, 1600), dark=i in (1, 6))
    proc.append(out)


def bg(slide, color=CREAM):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), SH)
    side.fill.solid(); side.fill.fore_color.rgb = GREEN
    side.line.fill.background()


def tbox(slide, l, t, w, h, text, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = 'Aptos'; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box


def tag(slide, l, t, w, text, fill=DEEP, tc=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Inches(0.3))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background()
    tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = 'Aptos'; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = tc
    return shp


def page(slide, n, dark=False):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.57), Inches(0.45), Inches(0.36), Inches(0.36))
    c.fill.solid(); c.fill.fore_color.rgb = WHITE if dark else GREEN
    c.line.fill.background()
    tf = c.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.name='Aptos'; r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb = DEEP if dark else WHITE


def footer_ref(slide, dark=False):
    col = WHITE if dark else MUTED
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(9.58), Inches(7.15), Inches(0.015))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(90,90,90) if dark else LINE
    line.line.fill.background()
    tbox(slide, Inches(0.55), Inches(9.72), Inches(2.8), Inches(0.2), '@dradaniely.freitas', size=9, color=col)
    # arrow + audio inspired marks
    arr = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(7.1), Inches(9.67), Inches(0.18), Inches(0.12))
    arr.fill.solid(); arr.fill.fore_color.rgb = col; arr.line.fill.background()
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.38), Inches(9.67), Inches(0.11), Inches(0.11))
    circ.fill.solid(); circ.fill.fore_color.rgb = col; circ.line.fill.background()
    ring = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(7.47), Inches(9.61), Inches(0.18), Inches(0.19))
    ring.fill.background(); ring.line.color.rgb = col


def card(slide, l, t, w, h, title, value, note):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = LINE
    tbox(slide, l + Inches(0.16), t + Inches(0.12), w - Inches(0.32), Inches(0.22), title.upper(), size=9, color=MUTED)
    tbox(slide, l + Inches(0.16), t + Inches(0.34), w - Inches(0.32), Inches(0.42), value, size=17, color=GREEN, bold=True)
    tbox(slide, l + Inches(0.16), t + Inches(0.77), w - Inches(0.32), Inches(0.26), note, size=8.5, color=TEXT)

# cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_picture(str(proc[0]), 0, 0, width=SW, height=SH)
over = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
over.fill.solid(); over.fill.fore_color.rgb = RGBColor(8, 12, 11); over.fill.transparency = 38; over.line.fill.background()
slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.72), Inches(0.14), Inches(8.9)).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = GREEN; slide.shapes[-1].line.fill.background()
tag(slide, Inches(0.72), Inches(0.72), Inches(2.25), 'DRA. DANIELY FREITAS', fill=WHITE, tc=DEEP)
tag(slide, Inches(3.03), Inches(0.72), Inches(1.68), 'VITAL SLIM', fill=GREEN, tc=WHITE)
tbox(slide, Inches(0.72), Inches(2.0), Inches(4.9), Inches(1.7), 'GLYNAC\nNA ROTINA', size=30, color=WHITE, bold=True)
tbox(slide, Inches(0.74), Inches(3.95), Inches(4.6), Inches(0.9), 'Glutationa, metabolismo e composição corporal com leitura clínica elegante e estratégica.', size=14.5, color=WHITE)
tbox(slide, Inches(0.74), Inches(8.65), Inches(4.0), Inches(0.3), 'Conteúdo clínico • premium • aplicável', size=10, color=WHITE)
footer_ref(slide, dark=True)
page(slide, 1, dark=True)

# science slide helper

def science(slide_no, img, chip, title, subtitle, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    photo_x = Inches(5.22)
    slide.shapes.add_picture(str(proc[img]), photo_x, Inches(0.78), width=Inches(2.62), height=Inches(8.05))
    o = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, photo_x, Inches(0.78), Inches(2.62), Inches(8.05))
    o.fill.solid(); o.fill.fore_color.rgb = RGBColor(9, 16, 14); o.fill.transparency = 48; o.line.fill.background()
    tag(slide, Inches(0.56), Inches(0.65), Inches(1.7), chip, fill=DEEP, tc=WHITE)
    tbox(slide, Inches(0.56), Inches(1.08), Inches(4.2), Inches(0.75), title, size=23, color=DEEP, bold=True)
    tbox(slide, Inches(0.56), Inches(1.82), Inches(4.05), Inches(0.52), subtitle, size=10.5, color=MUTED)
    positions = [(0.56,2.55),(2.57,2.55),(0.56,4.02),(2.57,4.02)]
    for (x,y),(a,b,c) in zip(positions, metrics):
        card(slide, Inches(x), Inches(y), Inches(1.78), Inches(1.18), a, b, c)
    ref = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.56), Inches(5.82), Inches(3.8), Inches(0.48))
    ref.fill.solid(); ref.fill.fore_color.rgb = MOSS; ref.line.fill.background()
    tbox(slide, Inches(0.75), Inches(5.95), Inches(3.3), Inches(0.18), 'Baseado em literatura recente sobre GlyNAC e glutationa.', size=8.5, color=GREEN)
    footer_ref(slide)
    page(slide, slide_no)

science(2, 1, 'PROVA 01', 'GLUTATIONA\nEM ALTA', 'Quando o protocolo é bem indicado, o ganho bioquímico fica mais nítido.', [
    ('Síntese de GSH', '~120%', 'produção antioxidante ampliada'),
    ('GSH total', '~55%', 'melhor reserva celular'),
    ('Relação GSH/GSSG', '~125%', 'redox mais favorável'),
    ('Radicais reativos', '~50%', 'estresse oxidativo menor'),
])
science(3, 2, 'PROVA 02', 'METABOLISMO\nMAIS FLEXÍVEL', 'Os marcadores sugerem resposta metabólica mais eficiente e adaptativa.', [
    ('NEFA em jejum', '~50%', 'mobilização lipídica maior'),
    ('NEFA alimentada', '~70%', 'uso energético mais eficaz'),
    ('Carboidratos', '~20%', 'melhor oxidação pós-refeição'),
    ('Supressão de NEFA', '~135%', 'sinal insulínico mais responsivo'),
])
science(4, 3, 'PROVA 03', 'COMPOSIÇÃO E\nFORÇA EM 14 DIAS', 'Mudanças elegantes são aquelas que respeitam o corpo e preservam performance.', [
    ('Gordura total', '21,6 → 20,0', 'queda objetiva de massa gorda'),
    ('Massa magra', '60,1 → 61,0', 'preservação com leve ganho'),
    ('Força dominante', '35 → 37', 'melhora funcional percebida'),
    ('Força não dominante', '31 → 34', 'ganho consistente bilateral'),
])

# practice
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
slide.shapes.add_picture(str(proc[4]), Inches(5.18), Inches(0.78), width=Inches(2.66), height=Inches(8.05))
o = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.18), Inches(0.78), Inches(2.66), Inches(8.05))
o.fill.solid(); o.fill.fore_color.rgb = RGBColor(9,16,14); o.fill.transparency = 52; o.line.fill.background()
tag(slide, Inches(0.56), Inches(0.65), Inches(1.92), 'MINHA PRÁTICA', fill=GREEN, tc=WHITE)
tbox(slide, Inches(0.56), Inches(1.08), Inches(4.15), Inches(0.82), 'PROTOCOLO COM\nCRITÉRIO CLÍNICO', size=22, color=DEEP, bold=True)
tbox(slide, Inches(0.56), Inches(1.9), Inches(4.0), Inches(0.48), 'Nada de modismo: avaliação, monitoramento e individualização.', size=10.5, color=MUTED)
panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.56), Inches(2.55), Inches(4.1), Inches(4.75))
panel.fill.solid(); panel.fill.fore_color.rgb = WHITE; panel.line.color.rgb = LINE
items = [
    'Avaliação de glutationa e contexto metabólico',
    'Homocisteína e função hepática monitoradas',
    'NAC e glicina fracionados conforme tolerância',
    'Insulina e composição corporal acompanhadas',
    'Ajustes integrados com TRH e nutrição personalizada',
]
y = 2.92
for item in items:
    b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.84), Inches(y+0.03), Inches(0.1), Inches(0.1))
    b.fill.solid(); b.fill.fore_color.rgb = GREEN; b.line.fill.background()
    tbox(slide, Inches(1.02), Inches(y), Inches(3.2), Inches(0.34), item, size=10.6, color=TEXT)
    y += 0.72
quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(6.82), Inches(3.55), Inches(0.42))
quote.fill.solid(); quote.fill.fore_color.rgb = MOSS; quote.line.fill.background()
tbox(slide, Inches(0.98), Inches(6.93), Inches(3.1), Inches(0.18), 'Sofisticação sem perder segurança clínica.', size=8.8, color=GREEN)
footer_ref(slide)
page(slide, 5)

# cta
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_picture(str(proc[5]), 0, 0, width=SW, height=SH)
over = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
over.fill.solid(); over.fill.fore_color.rgb = RGBColor(8, 12, 11); over.fill.transparency = 45; over.line.fill.background()
slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.74), Inches(0.14), Inches(8.8)).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = GREEN; slide.shapes[-1].line.fill.background()
tag(slide, Inches(0.72), Inches(0.72), Inches(2.1), 'ATENDIMENTO VITAL SLIM', fill=WHITE, tc=DEEP)
tbox(slide, Inches(0.72), Inches(2.0), Inches(5.0), Inches(1.5), 'NÃO É SOBRE\nTOMAR MAIS.', size=28, color=WHITE, bold=True)
tbox(slide, Inches(0.74), Inches(3.92), Inches(5.0), Inches(0.85), 'É sobre indicar melhor, monitorar de perto e ajustar com estratégia.', size=14.5, color=WHITE)
cta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.74), Inches(5.42), Inches(3.18), Inches(0.58))
cta.fill.solid(); cta.fill.fore_color.rgb = GREEN; cta.line.fill.background()
tf = cta.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'Agende sua avaliação integral'; r.font.name='Aptos'; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=WHITE
tbox(slide, Inches(0.74), Inches(8.7), Inches(4.2), Inches(0.26), '@dradaniely.freitas', size=10, color=WHITE)
footer_ref(slide, dark=True)
page(slide, 6, dark=True)

out = OUT / 'glynac-dra-daniely-vitalslim-v2-from-scratch.pptx'
prs.save(out)
print(out)

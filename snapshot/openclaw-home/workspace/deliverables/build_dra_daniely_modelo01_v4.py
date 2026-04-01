from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFilter, ImageEnhance, ImageOps
from pathlib import Path

OUT = Path('/root/.openclaw/workspace/deliverables')
TMP = OUT / 'tmp_modelo01_v4'
TMP.mkdir(parents=True, exist_ok=True)

SW, SH = 1000, 1250
GOLD = (159,136,68)
WHITE = RGBColor(255,255,255)
BLACK = RGBColor(18,18,18)
BROWN = RGBColor(87,71,23)
MUTED = RGBColor(100,100,100)
LIGHT_BG = RGBColor(250,248,244)
LIGHT_LINE = RGBColor(225,220,210)
GREY_BADGE = RGBColor(235,235,235)

photos = {
    '248': Path('/root/.openclaw/media/inbound/file_248---754c4855-061a-4e03-ae26-08b4a16b1da4.jpg'),
    '249': Path('/root/.openclaw/media/inbound/file_249---17318fde-dda9-4335-80a3-c22b8a774885.jpg'),
    '254': Path('/root/.openclaw/media/inbound/file_254---dfb15256-8d77-41ea-8aab-79b84f76357b.jpg'),
    '255': Path('/root/.openclaw/media/inbound/file_255---28c713b0-b988-4762-a1fe-0b96118a0f87.jpg'),
    '256': Path('/root/.openclaw/media/inbound/file_256---1133538b-77c4-4ae8-99a8-d8dedf3ccfb6.jpg'),
    '257': Path('/root/.openclaw/media/inbound/file_257---f0a37b69-9ccd-44c0-a981-4483bd299d9d.jpg'),
    '271': Path('/root/.openclaw/media/inbound/file_271---3696783f-b5d9-46dc-ba4d-3aedabbbbf0c.jpg'),
    '273': Path('/root/.openclaw/media/inbound/file_273---3f0ce689-e3db-48b7-85ca-971215f608f1.jpg'),
    '274': Path('/root/.openclaw/media/inbound/file_274---87a94e30-217c-4ff8-bb04-c012647be574.jpg'),
}


def crop_fill(img, size=(SW, SH)):
    tw, th = size
    src_ratio = img.width / img.height
    dst_ratio = tw / th
    if src_ratio > dst_ratio:
        nw = int(img.height * dst_ratio)
        left = (img.width - nw) // 2
        img = img.crop((left, 0, left + nw, img.height))
    else:
        nh = int(img.width / dst_ratio)
        top = (img.height - nh) // 2
        img = img.crop((0, top, img.width, top + nh))
    return img.resize((tw, th), Image.Resampling.LANCZOS)


def compose_editorial(photo_key, out_name, box, bg_crop='fill'):
    src = Image.open(photos[photo_key]).convert('RGB')
    base = crop_fill(src.copy())
    base = base.filter(ImageFilter.GaussianBlur(16))
    base = ImageEnhance.Brightness(base).enhance(0.32)
    base = ImageEnhance.Contrast(base).enhance(0.85)

    # foreground fully preserved inside a safe box (no destructive crop)
    x, y, w, h = box
    fg = src.copy()
    fg = ImageOps.contain(fg, (w, h), Image.Resampling.LANCZOS)

    canvas = base.convert('RGBA')
    fg_rgba = fg.convert('RGBA')
    px = x + (w - fg.width)//2
    py = y + (h - fg.height)//2
    canvas.alpha_composite(fg_rgba, (px, py))

    # dark editorial overlays flattened into the image
    overlay = Image.new('RGBA', (SW, SH), (0,0,0,0))
    d = ImageDraw.Draw(overlay)
    for y0 in range(SH):
        # stronger on bottom third for text
        if y0 < SH*0.35:
            alpha = int(18 + 25*(y0/(SH*0.35)))
        elif y0 < SH*0.65:
            alpha = int(45 + 40*((y0-SH*0.35)/(SH*0.30)))
        else:
            alpha = int(85 + 85*((y0-SH*0.65)/(SH*0.35)))
        d.line((0,y0,SW,y0), fill=(0,0,0,min(alpha,170)))
    canvas = Image.alpha_composite(canvas, overlay)

    out = canvas.convert('RGB')
    path = TMP / out_name
    out.save(path, quality=95)
    return path


def light_bg(out_name):
    img = Image.new('RGB', (SW,SH), (250,248,244))
    path = TMP / out_name
    img.save(path, quality=95)
    return path


def circle_avatar(photo_key, out_name):
    src = Image.open(photos[photo_key]).convert('RGB')
    fg = ImageOps.fit(src, (320,320), Image.Resampling.LANCZOS)
    mask = Image.new('L', (320,320), 0)
    ImageDraw.Draw(mask).ellipse((0,0,319,319), fill=255)
    out = Image.new('RGBA', (320,320), (0,0,0,0))
    out.paste(fg, (0,0), mask)
    path = TMP / out_name
    out.save(path)
    return path


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(8.333)
    prs.slide_height = Inches(10.417)
    return prs


def tbox(slide, l,t,w,h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l,t,w,h)
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name='Aptos'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return box


def pill(slide, l,t,w, txt, fill, tc):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l,t,w,Inches(0.29))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background()
    tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.name='Aptos'; r.font.size=Pt(9.5); r.font.bold=True; r.font.color.rgb=tc


def badge(slide, n):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.54), Inches(0.46), Inches(0.38), Inches(0.38))
    c.fill.solid(); c.fill.fore_color.rgb = GREY_BADGE; c.line.fill.background()
    tf = c.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.name='Aptos'; r.font.size=Pt(10.5); r.font.bold=True; r.font.color.rgb=BLACK


def dark_frame(slide):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.74), Inches(0.08), Inches(8.72))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(*GOLD); ln.line.fill.background()
    fl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(9.26), Inches(6.67), Inches(0.012))
    fl.fill.solid(); fl.fill.fore_color.rgb = WHITE; fl.line.fill.background()
    tbox(slide, Inches(0.63), Inches(9.34), Inches(1.5), Inches(0.16), '@dradaniely.freitas', size=8.2, color=WHITE)
    arr = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(6.87), Inches(9.32), Inches(0.16), Inches(0.10))
    arr.fill.solid(); arr.fill.fore_color.rgb = WHITE; arr.line.fill.background()


def light_frame(slide):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.74), Inches(0.08), Inches(8.72))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(*GOLD); ln.line.fill.background()
    fl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(9.26), Inches(6.67), Inches(0.012))
    fl.fill.solid(); fl.fill.fore_color.rgb = LIGHT_LINE; fl.line.fill.background()
    tbox(slide, Inches(0.63), Inches(9.34), Inches(1.5), Inches(0.16), '@dradaniely.freitas', size=8.2, color=MUTED)
    arr = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(6.87), Inches(9.32), Inches(0.16), Inches(0.10))
    arr.fill.solid(); arr.fill.fore_color.rgb = RGBColor(*GOLD); arr.line.fill.background()


def bgpic(slide, path):
    slide.shapes.add_picture(str(path), 0,0, width=Inches(8.333), height=Inches(10.417))


def cover(prs, bg, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bgpic(s, bg); dark_frame(s); badge(s,1)
    pill(s, Inches(0.64), Inches(0.76), Inches(2.05), 'DRA. DANIELY FREITAS', WHITE, BROWN)
    pill(s, Inches(2.78), Inches(0.76), Inches(1.42), 'VITAL SLIM', RGBColor(*GOLD), WHITE)
    tbox(s, Inches(0.72), Inches(5.18), Inches(4.7), Inches(1.4), title, size=23.5, color=WHITE, bold=True)
    tbox(s, Inches(0.72), Inches(6.68), Inches(5.4), Inches(0.7), subtitle, size=11.4, color=WHITE)
    return s


def proof(prs, bg, n, title, lead, body, ref):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bgpic(s, bg); dark_frame(s); badge(s,n)
    pill(s, Inches(0.64), Inches(0.76), Inches(2.05), 'DRA. DANIELY FREITAS', WHITE, BROWN)
    pill(s, Inches(2.78), Inches(0.76), Inches(1.42), 'VITAL SLIM', RGBColor(*GOLD), WHITE)
    tbox(s, Inches(0.72), Inches(5.10), Inches(5.0), Inches(1.0), title, size=21.5, color=WHITE, bold=True)
    tbox(s, Inches(0.72), Inches(6.22), Inches(5.6), Inches(0.48), lead, size=11.1, color=WHITE, bold=True)
    tbox(s, Inches(0.72), Inches(7.05), Inches(5.6), Inches(0.7), body, size=9.9, color=WHITE)
    tbox(s, Inches(0.72), Inches(8.82), Inches(5.8), Inches(0.18), ref, size=7.6, color=WHITE)
    return s


def practice(prs, avatar, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = light_bg('practice.jpg')
    bgpic(s,bg); light_frame(s); badge(s,5)
    s.shapes.add_picture(str(avatar), Inches(0.78), Inches(0.88), width=Inches(0.58), height=Inches(0.58))
    tbox(s, Inches(1.46), Inches(0.98), Inches(2.3), Inches(0.18), 'DRA. DANIELY FREITAS', size=8.8, color=BLACK, bold=True)
    tbox(s, Inches(1.46), Inches(1.16), Inches(2.3), Inches(0.18), '@dradaniely.freitas', size=8.4, color=MUTED)
    tbox(s, Inches(0.78), Inches(2.24), Inches(3.5), Inches(0.46), 'Minha prática:', size=22, color=RGBColor(*GOLD), bold=True)
    y = 3.10
    for b in bullets:
        tbox(s, Inches(0.82), Inches(y), Inches(6.2), Inches(0.42), f'• {b}', size=13.5, color=BLACK)
        y += 0.78
    return s


def cta(prs, bg, title, body, cta_text):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bgpic(s, bg); dark_frame(s); badge(s,6)
    pill(s, Inches(0.64), Inches(0.76), Inches(2.15), 'ATENDIMENTO VITAL SLIM', WHITE, BROWN)
    tbox(s, Inches(0.72), Inches(5.18), Inches(5.0), Inches(1.1), title, size=22.2, color=WHITE, bold=True)
    tbox(s, Inches(0.72), Inches(6.56), Inches(5.2), Inches(0.6), body, size=11.2, color=WHITE)
    tbox(s, Inches(0.72), Inches(7.72), Inches(3.8), Inches(0.36), cta_text, size=18.2, color=WHITE, bold=True)
    tbox(s, Inches(0.72), Inches(8.18), Inches(5.2), Inches(0.36), 'Se fizer sentido, me chame para uma avaliação estratégica.', size=9.2, color=WHITE)
    return s

# Compose safer editorial backgrounds with preserved body framing
cfg = {
    'g_cover': ('257', (300, 40, 620, 1120)),
    'g2': ('254', (310, 80, 610, 1060)),
    'g3': ('255', (300, 90, 620, 1040)),
    'g4': ('256', (300, 80, 620, 1060)),
    'g6': ('274', (280, 60, 640, 1080)),
    'h_cover': ('273', (270, 40, 650, 1120)),
    'h2': ('271', (310, 40, 600, 1110)),
    'h3': ('273', (270, 40, 650, 1120)),
    'h4': ('248', (340, 90, 560, 1050)),
    'h6': ('248', (340, 90, 560, 1050)),
}

BGS = {}
for key, (photo, box) in cfg.items():
    BGS[key] = compose_editorial(photo, f'{key}.jpg', box)
AV = circle_avatar('274', 'avatar.png')

# GLYNAC V4
prs = new_prs()
cover(prs, BGS['g_cover'], 'GLYNAC\nNA ROTINA', 'Aplicação clínica com mais clareza, critério e enquadramento fotográfico correto.')
proof(prs, BGS['g2'], 2, 'GLUTATIONA\nEM ALTA', 'Quando bem indicado, o protocolo pode favorecer um ambiente redox mais eficiente.', 'O ponto não é tendência. É contexto clínico, dose adequada e acompanhamento ao longo da resposta.', 'Referência editorial: Lai et al., 2024.')
proof(prs, BGS['g3'], 3, 'METABOLISMO\nMAIS FLEXÍVEL', 'A leitura clínica sugere uma resposta energética mais ajustada e estratégica.', 'Na prática, isso importa quando o objetivo é combinar performance, composição corporal e segurança metabólica.', 'Resumo visual para educação; individualização é indispensável.')
proof(prs, BGS['g4'], 4, 'COMPOSIÇÃO E\nFORÇA EM 14 DIAS', 'O alvo não é apenas reduzir gordura — é preservar função, massa e performance.', 'Em protocolos curtos, o que interessa é a direção clínica correta: menos ruído, mais critério e melhor execução.', 'Peça educativa; resultados dependem do caso.')
practice(prs, AV, [
    'Avalio contexto metabólico, inflamatório e objetivo clínico',
    'Defino se há indicação real antes de incluir qualquer protocolo',
    'Ajusto dose, timing e associação com nutrição',
    'Monitoro sinais, sintomas e exames ao longo do processo',
    'Reavalio o que está funcionando e o que precisa ser refinado',
])
cta(prs, BGS['g6'], 'NÃO É SOBRE\nTOMAR MAIS.', 'É sobre indicar melhor, acompanhar de perto e ajustar com estratégia.', 'Salve este post.')
path_g = OUT / '2026-04-01-glynac-modelo01-v4.pptx'
prs.save(path_g)
print(path_g)

# CAPILAR V4
prs = new_prs()
cover(prs, BGS['h_cover'], 'EMAGRECER SEM\nPERDER CABELO?', 'Sim — quando o processo respeita nutrientes, hormônios e velocidade de resposta.')
proof(prs, BGS['h2'], 2, 'A QUEDA NÃO\nSURGE DO NADA', 'Proteína baixa, ferritina, estresse e velocidade excessiva podem cobrar no fio.', 'Quando o emagrecimento vira agressão metabólica, o cabelo costuma ser um dos primeiros lugares a mostrar o custo biológico.', 'Referência editorial: eflúvio telógeno e estresse metabólico.')
proof(prs, BGS['h3'], 3, 'NO VITAL SLIM\nOLHAMOS A CAUSA', 'O foco não é esconder a queda. É corrigir o terreno que levou a ela.', 'Isso passa por exames, ajuste nutricional, revisão de estratégia e velocidade adequada para o seu caso.', 'Educação médica com abordagem integrada.')
proof(prs, BGS['h4'], 4, 'O OBJETIVO É\nEMAGRECER BEM', 'Resultado bonito é aquele que melhora corpo, energia e autoestima ao mesmo tempo.', 'Emagrecer bem é preservar massa, micronutrientes e confiança — sem transformar o processo em uma conta cara para o cabelo.', 'Peça editorial com linguagem acessível.')
practice(prs, AV, [
    'Mapeio exames e história clínica antes de acelerar o processo',
    'Corrijo proteína, ferro, vitaminas e sinais de carência',
    'Revejo a estratégia para emagrecer sem sacrificar saúde capilar',
    'Monitoro sintomas, composição corporal e resposta do fio',
    'Personalizo o ritmo de acordo com cada paciente',
])
cta(prs, BGS['h6'], 'EMAGRECER E\nPRESERVAR OS FIOS.', 'É possível quando existe diagnóstico, estratégia e acompanhamento.', 'Salve este post.')
path_h = OUT / '2026-04-01-capilar-modelo01-v4.pptx'
prs.save(path_h)
print(path_h)

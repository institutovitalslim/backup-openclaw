from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFilter, ImageEnhance, ImageOps
from pathlib import Path

OUT = Path('/root/.openclaw/workspace/deliverables')
TMP = OUT / 'tmp_modelo01_v5'
TMP.mkdir(parents=True, exist_ok=True)

SW, SH = 1000, 1250
GOLD = (159,136,68)
WHITE = RGBColor(255,255,255)
BLACK = RGBColor(20,20,20)
GREY = RGBColor(235,235,235)
MUTED = RGBColor(98,98,98)
LIGHT = RGBColor(249,247,243)
LINE = RGBColor(226,220,210)

photos = {
    '248': Path('/root/.openclaw/media/inbound/file_248---754c4855-061a-4e03-ae26-08b4a16b1da4.jpg'),
    '254': Path('/root/.openclaw/media/inbound/file_254---dfb15256-8d77-41ea-8aab-79b84f76357b.jpg'),
    '255': Path('/root/.openclaw/media/inbound/file_255---28c713b0-b988-4762-a1fe-0b96118a0f87.jpg'),
    '256': Path('/root/.openclaw/media/inbound/file_256---1133538b-77c4-4ae8-99a8-d8dedf3ccfb6.jpg'),
    '257': Path('/root/.openclaw/media/inbound/file_257---f0a37b69-9ccd-44c0-a981-4483bd299d9d.jpg'),
    '271': Path('/root/.openclaw/media/inbound/file_271---3696783f-b5d9-46dc-ba4d-3aedabbbbf0c.jpg'),
    '273': Path('/root/.openclaw/media/inbound/file_273---3f0ce689-e3db-48b7-85ca-971215f608f1.jpg'),
    '274': Path('/root/.openclaw/media/inbound/file_274---87a94e30-217c-4ff8-bb04-c012647be574.jpg'),
    '290': Path('/root/.openclaw/media/inbound/file_290---35d1e0bb-b0a3-4812-b210-92557bceaa5d.jpg'),
}


def crop_fill(img, size=(SW, SH)):
    tw, th = size
    src_ratio = img.width / img.height
    dst_ratio = tw / th
    if src_ratio > dst_ratio:
        nw = int(img.height * dst_ratio)
        left = (img.width - nw) // 2
        img = img.crop((left, 0, left + nw, img.height))
    else:
        nh = int(img.width / dst_ratio)
        top = (img.height - nh) // 2
        img = img.crop((0, top, img.width, top + nh))
    return img.resize((tw, th), Image.Resampling.LANCZOS)


def make_editorial(photo_key, name, box, blur=14):
    src = Image.open(photos[photo_key]).convert('RGB')
    bg = crop_fill(src.copy())
    bg = bg.filter(ImageFilter.GaussianBlur(blur))
    bg = ImageEnhance.Brightness(bg).enhance(0.28)
    bg = ImageEnhance.Contrast(bg).enhance(0.85)

    fg = ImageOps.contain(src.copy(), (box[2], box[3]), Image.Resampling.LANCZOS)
    canvas = bg.convert('RGBA')
    x = box[0] + (box[2] - fg.width)//2
    y = box[1] + (box[3] - fg.height)//2
    canvas.alpha_composite(fg.convert('RGBA'), (x, y))

    # flatten bottom dark band for text
    ov = Image.new('RGBA', (SW, SH), (0,0,0,0))
    d = ImageDraw.Draw(ov)
    for yy in range(SH):
        if yy < SH*0.48:
            alpha = int(18 + 20*(yy/(SH*0.48)))
        else:
            alpha = int(55 + 105*((yy-SH*0.48)/(SH*0.52)))
        d.line((0,yy,SW,yy), fill=(0,0,0,min(alpha,175)))
    # left gradient for readability
    for xx in range(SW):
        alpha = int(25 * (1 - min(xx, 500)/500)) if xx < 500 else 0
        d.line((xx,0,xx,SH), fill=(0,0,0,alpha))
    out = Image.alpha_composite(canvas, ov).convert('RGB')
    path = TMP / name
    out.save(path, quality=95)
    return path


def make_light(name):
    path = TMP / name
    Image.new('RGB', (SW,SH), (249,247,243)).save(path, quality=95)
    return path


def make_circle(photo_key, name):
    src = ImageOps.fit(Image.open(photos[photo_key]).convert('RGB'), (320,320), Image.Resampling.LANCZOS)
    mask = Image.new('L', (320,320), 0)
    ImageDraw.Draw(mask).ellipse((0,0,319,319), fill=255)
    out = Image.new('RGBA', (320,320), (0,0,0,0))
    out.paste(src, (0,0), mask)
    path = TMP / name
    out.save(path)
    return path


def prs_new():
    prs = Presentation()
    prs.slide_width = Inches(8.333)
    prs.slide_height = Inches(10.417)
    return prs


def box(slide, l,t,w,h, txt, size=18, color=WHITE, bold=False):
    tb = slide.shapes.add_textbox(l,t,w,h)
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = txt
    r.font.name = 'Aptos'; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb


def bgpic(slide, path):
    slide.shapes.add_picture(str(path), 0,0, width=Inches(8.333), height=Inches(10.417))


def top_identity(slide, dark=True):
    color = WHITE if dark else BLACK
    box(slide, Inches(0.66), Inches(0.78), Inches(2.6), Inches(0.18), 'DRA. DANIELY FREITAS', size=8.8, color=color, bold=True)
    box(slide, Inches(0.66), Inches(0.97), Inches(2.8), Inches(0.18), 'Instituto Vital Slim', size=8.4, color=color)


def left_bar(slide):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.74), Inches(0.08), Inches(8.72))
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor(*GOLD); r.line.fill.background()


def footer(slide, dark=True):
    c = WHITE if dark else LINE
    tc = WHITE if dark else MUTED
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(9.26), Inches(6.67), Inches(0.012))
    line.fill.solid(); line.fill.fore_color.rgb = c; line.line.fill.background()
    box(slide, Inches(0.63), Inches(9.34), Inches(1.7), Inches(0.16), '@dradaniely.freitas', size=8.2, color=tc)
    arr = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(6.87), Inches(9.32), Inches(0.16), Inches(0.10))
    arr.fill.solid(); arr.fill.fore_color.rgb = tc; arr.line.fill.background()


def badge(slide, n):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.54), Inches(0.46), Inches(0.38), Inches(0.38))
    c.fill.solid(); c.fill.fore_color.rgb = GREY; c.line.fill.background()
    tf = c.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.name='Aptos'; r.font.size=Pt(10.5); r.font.bold=True; r.font.color.rgb=BLACK


def dark_slide(prs, bg, n, title, sub, body=None, ref=None, cta=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bgpic(s,bg); left_bar(s); top_identity(s, dark=True); footer(s, dark=True); badge(s,n)
    box(s, Inches(0.72), Inches(5.10), Inches(5.0), Inches(1.2), title, size=22.5, color=WHITE, bold=True)
    box(s, Inches(0.72), Inches(6.44), Inches(5.25), Inches(0.55), sub, size=11.5, color=WHITE, bold=True)
    if body:
        box(s, Inches(0.72), Inches(7.10), Inches(5.35), Inches(0.75), body, size=9.7, color=WHITE)
    if ref:
        box(s, Inches(0.72), Inches(8.82), Inches(5.5), Inches(0.18), ref, size=7.4, color=WHITE)
    if cta:
        box(s, Inches(0.72), Inches(7.80), Inches(4.5), Inches(0.40), cta, size=18.2, color=WHITE, bold=True)
        box(s, Inches(0.72), Inches(8.20), Inches(5.3), Inches(0.36), 'Se fizer sentido, me chame para uma avaliação estratégica.', size=9.1, color=WHITE)
    return s


def practice_slide(prs, avatar, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bgpic(s, make_light('light.jpg'))
    left_bar(s); footer(s, dark=False); badge(s,5)
    s.shapes.add_picture(str(avatar), Inches(0.78), Inches(0.88), width=Inches(0.58), height=Inches(0.58))
    box(s, Inches(1.46), Inches(0.98), Inches(2.8), Inches(0.18), 'DRA. DANIELY FREITAS', size=8.8, color=BLACK, bold=True)
    box(s, Inches(1.46), Inches(1.16), Inches(2.8), Inches(0.18), '@dradaniely.freitas', size=8.4, color=MUTED)
    box(s, Inches(0.78), Inches(2.24), Inches(3.6), Inches(0.46), 'Minha prática:', size=22, color=RGBColor(*GOLD), bold=True)
    y = 3.12
    for b in bullets:
        box(s, Inches(0.82), Inches(y), Inches(6.25), Inches(0.42), f'• {b}', size=13.4, color=BLACK)
        y += 0.78
    return s

cfg = {
    'g1': ('274', (335, 50, 560, 1110)),
    'g2': ('254', (330, 90, 570, 1040)),
    'g3': ('255', (330, 90, 570, 1040)),
    'g4': ('256', (330, 90, 570, 1040)),
    'g6': ('257', (340, 60, 550, 1090)),
    'h1': ('273', (320, 40, 590, 1120)),
    'h2': ('271', (340, 50, 560, 1110)),
    'h3': ('274', (330, 50, 570, 1100)),
    'h4': ('248', (360, 110, 520, 1030)),
    'h6': ('273', (320, 40, 590, 1120)),
}
BG = {k: make_editorial(photo, f'{k}.jpg', box) for k,(photo,box) in cfg.items()}
AV = make_circle('290', 'avatar-290.png')

# GLYNAC V5
prs = prs_new()
dark_slide(prs, BG['g1'], 1, 'GLYNAC\nNA ROTINA', 'Aplicação clínica com mais clareza, critério e autoridade visual.')
dark_slide(prs, BG['g2'], 2, 'GLUTATIONA\nEM ALTA', 'Quando bem indicado, o protocolo pode favorecer um ambiente redox mais eficiente.', 'O ponto não é tendência. É contexto clínico, dose adequada e acompanhamento ao longo da resposta.', 'Referência editorial: Lai et al., 2024.')
dark_slide(prs, BG['g3'], 3, 'METABOLISMO\nMAIS FLEXÍVEL', 'A leitura clínica sugere uma resposta energética mais ajustada e estratégica.', 'Na prática, isso importa quando o objetivo é combinar performance, composição corporal e segurança metabólica.', 'Resumo visual para educação; individualização é indispensável.')
dark_slide(prs, BG['g4'], 4, 'COMPOSIÇÃO E\nFORÇA EM 14 DIAS', 'O alvo não é apenas reduzir gordura — é preservar função, massa e performance.', 'Em protocolos curtos, o que interessa é a direção clínica correta: menos ruído, mais critério e melhor execução.', 'Peça educativa; resultados dependem do caso.')
practice_slide(prs, AV, [
    'Avalio contexto metabólico, inflamatório e objetivo clínico',
    'Defino se há indicação real antes de incluir qualquer protocolo',
    'Ajusto dose, timing e associação com nutrição',
    'Monitoro sinais, sintomas e exames ao longo do processo',
    'Reavalio o que está funcionando e o que precisa ser refinado',
])
dark_slide(prs, BG['g6'], 6, 'NÃO É SOBRE\nTOMAR MAIS.', 'É sobre indicar melhor, acompanhar de perto e ajustar com estratégia.', cta='Salve este post.')
path_g = OUT / '2026-04-01-glynac-modelo01-v5.pptx'
prs.save(path_g)
print(path_g)

# CAPILAR V5
prs = prs_new()
dark_slide(prs, BG['h1'], 1, 'EMAGRECER SEM\nPERDER CABELO?', 'Sim — quando o processo respeita nutrientes, hormônios e velocidade de resposta.')
dark_slide(prs, BG['h2'], 2, 'A QUEDA NÃO\nSURGE DO NADA', 'Proteína baixa, ferritina, estresse e velocidade excessiva podem cobrar no fio.', 'Quando o emagrecimento vira agressão metabólica, o cabelo costuma ser um dos primeiros lugares a mostrar o custo biológico.', 'Referência editorial: eflúvio telógeno e estresse metabólico.')
dark_slide(prs, BG['h3'], 3, 'NO VITAL SLIM\nOLHAMOS A CAUSA', 'O foco não é esconder a queda. É corrigir o terreno que levou a ela.', 'Isso passa por exames, ajuste nutricional, revisão de estratégia e velocidade adequada para o seu caso.', 'Educação médica com abordagem integrada.')
dark_slide(prs, BG['h4'], 4, 'O OBJETIVO É\nEMAGRECER BEM', 'Resultado bonito é aquele que melhora corpo, energia e autoestima ao mesmo tempo.', 'Emagrecer bem é preservar massa, micronutrientes e confiança — sem transformar o processo em uma conta cara para o cabelo.', 'Peça editorial com linguagem acessível.')
practice_slide(prs, AV, [
    'Mapeio exames e história clínica antes de acelerar o processo',
    'Corrijo proteína, ferro, vitaminas e sinais de carência',
    'Revejo a estratégia para emagrecer sem sacrificar saúde capilar',
    'Monitoro sintomas, composição corporal e resposta do fio',
    'Personalizo o ritmo de acordo com cada paciente',
])
dark_slide(prs, BG['h6'], 6, 'EMAGRECER E\nPRESERVAR OS FIOS.', 'É possível quando existe diagnóstico, estratégia e acompanhamento.', cta='Salve este post.')
path_h = OUT / '2026-04-01-capilar-modelo01-v5.pptx'
prs.save(path_h)
print(path_h)

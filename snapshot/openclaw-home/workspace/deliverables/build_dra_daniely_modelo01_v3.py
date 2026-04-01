from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFilter
from pathlib import Path

OUT = Path('/root/.openclaw/workspace/deliverables')
TMP = OUT / 'tmp_modelo01_v3'
TMP.mkdir(parents=True, exist_ok=True)

SW, SH = 1000, 1250  # 4:5
GOLD = (159,136,68)
BROWN = RGBColor(87,71,23)
WHITE = RGBColor(255,255,255)
BLACK = RGBColor(20,20,20)
MUTED = RGBColor(100,100,100)
LIGHT_BG = RGBColor(250,248,244)
LIGHT_LINE = RGBColor(225,220,210)
GREY_BADGE = RGBColor(235,235,235)

photos = {
    '248': Path('/root/.openclaw/media/inbound/file_248---754c4855-061a-4e03-ae26-08b4a16b1da4.jpg'),
    '249': Path('/root/.openclaw/media/inbound/file_249---17318fde-dda9-4335-80a3-c22b8a774885.jpg'),
    '254': Path('/root/.openclaw/media/inbound/file_254---dfb15256-8d77-41ea-8aab-79b84f76357b.jpg'),
    '255': Path('/root/.openclaw/media/inbound/file_255---28c713b0-b988-4762-a1fe-0b96118a0f87.jpg'),
    '256': Path('/root/.openclaw/media/inbound/file_256---1133538b-77c4-4ae8-99a8-d8dedf3ccfb6.jpg'),
    '257': Path('/root/.openclaw/media/inbound/file_257---f0a37b69-9ccd-44c0-a981-4483bd299d9d.jpg'),
    '271': Path('/root/.openclaw/media/inbound/file_271---3696783f-b5d9-46dc-ba4d-3aedabbbbf0c.jpg'),
    '273': Path('/root/.openclaw/media/inbound/file_273---3f0ce689-e3db-48b7-85ca-971215f608f1.jpg'),
    '274': Path('/root/.openclaw/media/inbound/file_274---87a94e30-217c-4ff8-bb04-c012647be574.jpg'),
}


def crop_fill(path, size=(SW, SH), anchor='center'):
    img = Image.open(path).convert('RGB')
    tw, th = size
    src_ratio = img.width / img.height
    dst_ratio = tw / th
    if src_ratio > dst_ratio:
        nw = int(img.height * dst_ratio)
        if anchor == 'left':
            left = 0
        elif anchor == 'right':
            left = img.width - nw
        else:
            left = (img.width - nw) // 2
        img = img.crop((left, 0, left + nw, img.height))
    else:
        nh = int(img.width / dst_ratio)
        top = (img.height - nh) // 2
        img = img.crop((0, top, img.width, top + nh))
    return img.resize((tw, th), Image.Resampling.LANCZOS)


def dark_editorial_bg(photo_key, name, anchor='center'):
    img = crop_fill(photos[photo_key], anchor=anchor)
    overlay = Image.new('RGBA', (SW, SH), (0,0,0,0))
    d = ImageDraw.Draw(overlay)
    # Strong bottom overlay, softer top overlay, slight left vignette
    for y in range(SH):
        if y < SH*0.45:
            alpha = int(45 + 35*(y/(SH*0.45)))
        else:
            alpha = int(90 + 110*((y - SH*0.45)/(SH*0.55)))
        d.line((0, y, SW, y), fill=(0,0,0, min(alpha,190)))
    vignette = Image.new('RGBA', (SW, SH), (0,0,0,0))
    dv = ImageDraw.Draw(vignette)
    for x in range(SW):
        alpha = int(45 * (1 - x/SW))
        dv.line((x, 0, x, SH), fill=(0,0,0,alpha))
    out = Image.alpha_composite(img.convert('RGBA'), overlay)
    out = Image.alpha_composite(out, vignette).convert('RGB')
    path = TMP / name
    out.save(path, quality=95)
    return path


def light_bg(name):
    img = Image.new('RGB', (SW, SH), (250,248,244))
    path = TMP / name
    img.save(path, quality=95)
    return path


def circle_avatar(photo_key, name):
    img = crop_fill(photos[photo_key], size=(320,320))
    mask = Image.new('L', (320,320), 0)
    d = ImageDraw.Draw(mask)
    d.ellipse((0,0,319,319), fill=255)
    out = Image.new('RGBA', (320,320), (0,0,0,0))
    out.paste(img, (0,0), mask)
    path = TMP / name
    out.save(path)
    return path


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(8.333)
    prs.slide_height = Inches(10.417)
    return prs


def textbox(slide, l,t,w,h, txt, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l,t,w,h)
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.name = 'Aptos'; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box


def pill(slide, l,t,w, txt, fill_rgb, text_rgb):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l,t,w,Inches(0.29))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill_rgb; shp.line.fill.background()
    tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.name = 'Aptos'; r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = text_rgb


def badge(slide, n):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.54), Inches(0.46), Inches(0.38), Inches(0.38))
    c.fill.solid(); c.fill.fore_color.rgb = GREY_BADGE; c.line.fill.background()
    tf = c.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.name = 'Aptos'; r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = BLACK


def frame(slide, dark=True):
    # left gold line
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.74), Inches(0.08), Inches(8.72))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(*GOLD); ln.line.fill.background()
    # footer line
    fl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(9.26), Inches(6.67), Inches(0.012))
    fl.fill.solid(); fl.fill.fore_color.rgb = RGBColor(255,255,255) if dark else LIGHT_LINE; fl.line.fill.background()
    # handle
    textbox(slide, Inches(0.63), Inches(9.34), Inches(1.5), Inches(0.16), '@dradaniely.freitas', size=8.2, color=WHITE if dark else MUTED)
    # arrow
    arr = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(6.87), Inches(9.32), Inches(0.16), Inches(0.10))
    arr.fill.solid(); arr.fill.fore_color.rgb = WHITE if dark else RGBColor(*GOLD); arr.line.fill.background()


def place_bg(slide, path):
    slide.shapes.add_picture(str(path), 0, 0, width=Inches(8.333), height=Inches(10.417))


def cover_slide(prs, bg_path, title, subtitle, page_no=1):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    place_bg(s, bg_path)
    frame(s, dark=True)
    pill(s, Inches(0.64), Inches(0.76), Inches(2.05), 'DRA. DANIELY FREITAS', WHITE, BROWN)
    pill(s, Inches(2.78), Inches(0.76), Inches(1.42), 'VITAL SLIM', RGBColor(*GOLD), WHITE)
    badge(s, page_no)
    textbox(s, Inches(0.72), Inches(5.10), Inches(5.25), Inches(1.45), title, size=24, color=WHITE, bold=True)
    textbox(s, Inches(0.72), Inches(6.68), Inches(5.3), Inches(0.65), subtitle, size=11.8, color=WHITE)
    return s


def evidence_slide(prs, bg_path, page_no, title, lead, body, ref):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    place_bg(s, bg_path)
    frame(s, dark=True)
    pill(s, Inches(0.64), Inches(0.76), Inches(2.05), 'DRA. DANIELY FREITAS', WHITE, BROWN)
    pill(s, Inches(2.78), Inches(0.76), Inches(1.42), 'VITAL SLIM', RGBColor(*GOLD), WHITE)
    badge(s, page_no)
    textbox(s, Inches(0.72), Inches(5.00), Inches(5.5), Inches(0.95), title, size=22.5, color=WHITE, bold=True)
    textbox(s, Inches(0.72), Inches(6.08), Inches(5.9), Inches(0.65), lead, size=12.2, color=WHITE, bold=True)
    textbox(s, Inches(0.72), Inches(6.78), Inches(5.9), Inches(0.86), body, size=10.6, color=WHITE)
    textbox(s, Inches(0.72), Inches(8.82), Inches(5.9), Inches(0.18), ref, size=7.8, color=WHITE)
    return s


def practice_slide(prs, avatar_path, title, bullets, page_no=5):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = light_bg(f'practice-{page_no}.jpg')
    place_bg(s, bg)
    # frame in light mode
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.74), Inches(0.08), Inches(8.72))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(*GOLD); ln.line.fill.background()
    fl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(9.26), Inches(6.67), Inches(0.012))
    fl.fill.solid(); fl.fill.fore_color.rgb = LIGHT_LINE; fl.line.fill.background()
    textbox(s, Inches(0.63), Inches(9.34), Inches(1.5), Inches(0.16), '@dradaniely.freitas', size=8.2, color=MUTED)
    arr = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(6.87), Inches(9.32), Inches(0.16), Inches(0.10))
    arr.fill.solid(); arr.fill.fore_color.rgb = RGBColor(*GOLD); arr.line.fill.background()
    badge(s, page_no)
    s.shapes.add_picture(str(avatar_path), Inches(0.78), Inches(0.88), width=Inches(0.58), height=Inches(0.58))
    textbox(s, Inches(1.46), Inches(0.98), Inches(2.3), Inches(0.18), 'DRA. DANIELY FREITAS', size=8.8, color=BLACK, bold=True)
    textbox(s, Inches(1.46), Inches(1.16), Inches(2.3), Inches(0.18), '@dradaniely.freitas', size=8.4, color=MUTED)
    textbox(s, Inches(0.78), Inches(2.25), Inches(3.6), Inches(0.5), title, size=22, color=RGBColor(*GOLD), bold=True)
    y = 3.12
    for b in bullets:
        textbox(s, Inches(0.82), Inches(y), Inches(6.0), Inches(0.44), f'• {b}', size=14.2, color=BLACK)
        y += 0.78
    return s


def cta_slide(prs, bg_path, title, body, cta, page_no=6):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    place_bg(s, bg_path)
    frame(s, dark=True)
    pill(s, Inches(0.64), Inches(0.76), Inches(2.15), 'ATENDIMENTO VITAL SLIM', WHITE, BROWN)
    badge(s, page_no)
    textbox(s, Inches(0.72), Inches(5.06), Inches(5.3), Inches(1.25), title, size=22.8, color=WHITE, bold=True)
    textbox(s, Inches(0.72), Inches(6.52), Inches(5.5), Inches(0.70), body, size=11.7, color=WHITE)
    textbox(s, Inches(0.72), Inches(7.72), Inches(4.2), Inches(0.42), cta, size=18.5, color=WHITE, bold=True)
    textbox(s, Inches(0.72), Inches(8.16), Inches(5.2), Inches(0.42), 'Se fizer sentido, me chame para uma avaliação estratégica.', size=9.5, color=WHITE)
    return s

# backgrounds with flattened overlays
BG = {
    'g_cover': dark_editorial_bg('257', 'g-cover.jpg', 'center'),
    'g_2': dark_editorial_bg('254', 'g-2.jpg', 'center'),
    'g_3': dark_editorial_bg('255', 'g-3.jpg', 'center'),
    'g_4': dark_editorial_bg('256', 'g-4.jpg', 'center'),
    'g_6': dark_editorial_bg('271', 'g-6.jpg', 'center'),
    'h_cover': dark_editorial_bg('271', 'h-cover.jpg', 'center'),
    'h_2': dark_editorial_bg('273', 'h-2.jpg', 'center'),
    'h_3': dark_editorial_bg('257', 'h-3.jpg', 'center'),
    'h_4': dark_editorial_bg('249', 'h-4.jpg', 'center'),
    'h_6': dark_editorial_bg('248', 'h-6.jpg', 'center'),
    'avatar': circle_avatar('274', 'avatar.png'),
}

# GLYNAC
prs = new_prs()
cover_slide(prs, BG['g_cover'], 'GLYNAC\nNA ROTINA', 'Aplicação clínica com mais clareza, critério e autoridade visual.')
evidence_slide(prs, BG['g_2'], 2, 'GLUTATIONA\nEM ALTA', 'Quando bem indicado, o protocolo pode favorecer um ambiente redox mais eficiente.', 'O ponto não é tendência. É contexto clínico, dose adequada e acompanhamento ao longo da resposta.', 'Referência editorial: Lai et al., 2024.')
evidence_slide(prs, BG['g_3'], 3, 'METABOLISMO\nMAIS FLEXÍVEL', 'A leitura clínica sugere uma resposta energética mais ajustada e estratégica.', 'Na prática, isso importa quando o objetivo é combinar performance, composição corporal e segurança metabólica.', 'Resumo visual para educação; individualização é indispensável.')
evidence_slide(prs, BG['g_4'], 4, 'COMPOSIÇÃO E\nFORÇA EM 14 DIAS', 'O alvo não é apenas reduzir gordura — é preservar função, massa e performance.', 'Em protocolos curtos, o que interessa é a direção clínica correta: menos ruído, mais critério e melhor execução.', 'Peça educativa; resultados dependem do caso.')
practice_slide(prs, BG['avatar'], 'Minha prática:', [
    'Avalio contexto metabólico, inflamatório e objetivo clínico',
    'Defino se há indicação real antes de incluir qualquer protocolo',
    'Ajusto dose, timing e associação com nutrição',
    'Monitoro sinais, sintomas e exames ao longo do processo',
    'Reavalio o que está funcionando e o que precisa ser refinado',
])
cta_slide(prs, BG['g_6'], 'NÃO É SOBRE\nTOMAR MAIS.', 'É sobre indicar melhor, acompanhar de perto e ajustar com estratégia.', 'Salve este post.')
path_g = OUT / '2026-04-01-glynac-modelo01-v3.pptx'
prs.save(path_g)
print(path_g)

# Hair loss
prs = new_prs()
cover_slide(prs, BG['h_cover'], 'EMAGRECER SEM\nPERDER CABELO?', 'Sim — quando o processo respeita nutrientes, hormônios e velocidade de resposta.')
evidence_slide(prs, BG['h_2'], 2, 'A QUEDA NÃO\nSURGE DO NADA', 'Proteína baixa, ferritina, estresse e velocidade excessiva podem cobrar no fio.', 'Quando o emagrecimento vira agressão metabólica, o cabelo costuma ser um dos primeiros lugares a mostrar o custo biológico.', 'Referência editorial: eflúvio telógeno e estresse metabólico.')
evidence_slide(prs, BG['h_3'], 3, 'NO VITAL SLIM\nOLHAMOS A CAUSA', 'O foco não é esconder a queda. É corrigir o terreno que levou a ela.', 'Isso passa por exames, ajuste nutricional, revisão de estratégia e velocidade adequada para o seu caso.', 'Educação médica com abordagem integrada.')
evidence_slide(prs, BG['h_4'], 4, 'O OBJETIVO É\nEMAGRECER BEM', 'Resultado bonito é aquele que melhora corpo, energia e autoestima ao mesmo tempo.', 'Emagrecer bem é preservar massa, micronutrientes e confiança — sem transformar o processo em uma conta cara para o cabelo.', 'Peça editorial com linguagem acessível.')
practice_slide(prs, BG['avatar'], 'Minha prática:', [
    'Mapeio exames e história clínica antes de acelerar o processo',
    'Corrijo proteína, ferro, vitaminas e sinais de carência',
    'Revejo a estratégia para emagrecer sem sacrificar saúde capilar',
    'Monitoro sintomas, composição corporal e resposta do fio',
    'Personalizo o ritmo de acordo com cada paciente',
])
cta_slide(prs, BG['h_6'], 'EMAGRECER E\nPRESERVAR OS FIOS.', 'É possível quando existe diagnóstico, estratégia e acompanhamento.', 'Salve este post.')
path_h = OUT / '2026-04-01-capilar-modelo01-v3.pptx'
prs.save(path_h)
print(path_h)

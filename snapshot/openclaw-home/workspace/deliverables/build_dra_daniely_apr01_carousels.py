from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
from pathlib import Path

OUT = Path('/root/.openclaw/workspace/deliverables')
TMP = OUT / 'tmp_apr01_carousels'
TMP.mkdir(parents=True, exist_ok=True)

# Brand palette from manual
GOLD = RGBColor(159, 136, 68)     # #9F8844
BROWN = RGBColor(87, 71, 23)      # #574717
CREAM = RGBColor(247, 243, 236)
CHARCOAL = RGBColor(22, 22, 22)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(116, 104, 83)
LINE = RGBColor(222, 216, 206)
BEIGE = RGBColor(235, 227, 214)

photos = {
    '248': Path('/root/.openclaw/media/inbound/file_248---754c4855-061a-4e03-ae26-08b4a16b1da4.jpg'),
    '249': Path('/root/.openclaw/media/inbound/file_249---17318fde-dda9-4335-80a3-c22b8a774885.jpg'),
    '254': Path('/root/.openclaw/media/inbound/file_254---dfb15256-8d77-41ea-8aab-79b84f76357b.jpg'),
    '255': Path('/root/.openclaw/media/inbound/file_255---28c713b0-b988-4762-a1fe-0b96118a0f87.jpg'),
    '256': Path('/root/.openclaw/media/inbound/file_256---1133538b-77c4-4ae8-99a8-d8dedf3ccfb6.jpg'),
    '257': Path('/root/.openclaw/media/inbound/file_257---f0a37b69-9ccd-44c0-a981-4483bd299d9d.jpg'),
    '270': Path('/root/.openclaw/media/inbound/file_270---f51ce3ef-c6d1-4371-9d26-841f6f863a79.jpg'),
    '271': Path('/root/.openclaw/media/inbound/file_271---3696783f-b5d9-46dc-ba4d-3aedabbbbf0c.jpg'),
    '273': Path('/root/.openclaw/media/inbound/file_273---3f0ce689-e3db-48b7-85ca-971215f608f1.jpg'),
    '274': Path('/root/.openclaw/media/inbound/file_274---87a94e30-217c-4ff8-bb04-c012647be574.jpg'),
}

def crop_fill(src, dst, size, dark=False):
    img = Image.open(src).convert('RGB')
    tw, th = size
    src_ratio = img.width / img.height
    dst_ratio = tw / th
    if src_ratio > dst_ratio:
        nw = int(img.height * dst_ratio)
        left = (img.width - nw) // 2
        img = img.crop((left, 0, left + nw, img.height))
    else:
        nh = int(img.width / dst_ratio)
        top = (img.height - nh) // 2
        img = img.crop((0, top, img.width, top + nh))
    img = img.resize((tw, th), Image.Resampling.LANCZOS)
    if dark:
        overlay = Image.new('RGBA', (tw, th), (0, 0, 0, 0))
        draw = ImageDraw.Draw(overlay)
        for y in range(th):
            alpha = int(140 * (y / th))
            draw.line((0, y, tw, y), fill=(15, 15, 15, alpha))
        for x in range(tw):
            alpha = int(90 * ((tw - x) / tw))
            draw.line((x, 0, x, th), fill=(15, 15, 15, alpha))
        img = Image.alpha_composite(img.convert('RGBA'), overlay).convert('RGB')
    img.save(dst, quality=95)
    return dst

proc = {}
for key, path in photos.items():
    out = TMP / f'{key}.jpg'
    crop_fill(path, out, (1200, 1600), dark=key in {'249','257','271'})
    proc[key] = out


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(8.333)
    prs.slide_height = Inches(10.417)
    return prs


def add_bg(slide, prs, color=CREAM):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), prs.slide_height)
    side.fill.solid(); side.fill.fore_color.rgb = GOLD; side.line.fill.background()


def tbox(slide, l, t, w, h, text, size=18, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = 'Aptos'; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box


def tag(slide, l, t, w, text, fill=BROWN, tc=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Inches(0.3))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background()
    tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = 'Aptos'; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = tc
    return shp


def page_badge(slide, n, dark=False):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.56), Inches(0.44), Inches(0.38), Inches(0.38))
    c.fill.solid(); c.fill.fore_color.rgb = WHITE if dark else GOLD; c.line.fill.background()
    tf = c.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.name = 'Aptos'; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = BROWN if dark else WHITE


def footer(slide, dark=False):
    col = WHITE if dark else MUTED
    line_col = RGBColor(110,110,110) if dark else LINE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(9.58), Inches(7.15), Inches(0.015))
    line.fill.solid(); line.fill.fore_color.rgb = line_col; line.line.fill.background()
    tbox(slide, Inches(0.56), Inches(9.72), Inches(2.8), Inches(0.2), '@dradaniely.freitas', size=9, color=col)
    arr = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(7.1), Inches(9.67), Inches(0.18), Inches(0.12))
    arr.fill.solid(); arr.fill.fore_color.rgb = col; arr.line.fill.background()


def card(slide, l, t, w, h, title, body):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE; shp.line.color.rgb = LINE
    tbox(slide, l + Inches(0.15), t + Inches(0.11), w - Inches(0.3), Inches(0.23), title.upper(), size=9, color=GOLD, bold=True)
    tbox(slide, l + Inches(0.15), t + Inches(0.33), w - Inches(0.3), h - Inches(0.42), body, size=10.2, color=CHARCOAL)


def metric_card(slide, l, t, w, h, title, value, note):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE; shp.line.color.rgb = LINE
    tbox(slide, l + Inches(0.16), t + Inches(0.12), w - Inches(0.32), Inches(0.18), title.upper(), size=8.8, color=MUTED)
    tbox(slide, l + Inches(0.16), t + Inches(0.31), w - Inches(0.32), Inches(0.34), value, size=16, color=BROWN, bold=True)
    tbox(slide, l + Inches(0.16), t + Inches(0.68), w - Inches(0.32), Inches(0.27), note, size=8.4, color=CHARCOAL)


def cover_slide(prs, photo_key, chip1, chip2, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(proc[photo_key]), 0, 0, width=prs.slide_width, height=prs.slide_height)
    over = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    over.fill.solid(); over.fill.fore_color.rgb = RGBColor(18,18,18); over.fill.transparency = 44; over.line.fill.background()
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.72), Inches(0.14), Inches(8.85))
    side.fill.solid(); side.fill.fore_color.rgb = GOLD; side.line.fill.background()
    tag(slide, Inches(0.72), Inches(0.72), Inches(2.2), chip1, fill=WHITE, tc=BROWN)
    tag(slide, Inches(2.98), Inches(0.72), Inches(1.55), chip2, fill=GOLD, tc=WHITE)
    tbox(slide, Inches(0.72), Inches(1.95), Inches(5.0), Inches(1.7), title, size=28, color=WHITE, bold=True)
    tbox(slide, Inches(0.74), Inches(3.9), Inches(4.8), Inches(0.95), subtitle, size=14.2, color=WHITE)
    footer(slide, dark=True)
    page_badge(slide, 1, dark=True)
    return slide


def science_slide(prs, photo_key, page_no, chip, title, subtitle, metrics, ref='Referência científica no rodapé para apoio editorial.'):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    slide.shapes.add_picture(str(proc[photo_key]), Inches(5.18), Inches(0.78), width=Inches(2.66), height=Inches(8.05))
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.18), Inches(0.78), Inches(2.66), Inches(8.05))
    overlay.fill.solid(); overlay.fill.fore_color.rgb = RGBColor(20,20,20); overlay.fill.transparency = 52; overlay.line.fill.background()
    tag(slide, Inches(0.56), Inches(0.64), Inches(1.58), chip, fill=BROWN, tc=WHITE)
    tbox(slide, Inches(0.56), Inches(1.04), Inches(4.2), Inches(0.8), title, size=22, color=BROWN, bold=True)
    tbox(slide, Inches(0.56), Inches(1.82), Inches(4.05), Inches(0.5), subtitle, size=10.2, color=MUTED)
    positions = [(0.56,2.55),(2.57,2.55),(0.56,4.0),(2.57,4.0)]
    for (x,y),(a,b,c) in zip(positions, metrics):
        metric_card(slide, Inches(x), Inches(y), Inches(1.78), Inches(1.14), a, b, c)
    refbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.56), Inches(5.8), Inches(3.9), Inches(0.46))
    refbox.fill.solid(); refbox.fill.fore_color.rgb = BEIGE; refbox.line.fill.background()
    tbox(slide, Inches(0.76), Inches(5.93), Inches(3.45), Inches(0.18), ref, size=8.1, color=BROWN)
    footer(slide)
    page_badge(slide, page_no)
    return slide


def practice_slide(prs, photo_key, title, subtitle, bullets, page_no=5):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    slide.shapes.add_picture(str(proc[photo_key]), Inches(5.18), Inches(0.78), width=Inches(2.66), height=Inches(8.05))
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.18), Inches(0.78), Inches(2.66), Inches(8.05))
    overlay.fill.solid(); overlay.fill.fore_color.rgb = RGBColor(20,20,20); overlay.fill.transparency = 55; overlay.line.fill.background()
    tag(slide, Inches(0.56), Inches(0.64), Inches(1.88), 'MINHA PRÁTICA', fill=GOLD, tc=WHITE)
    tbox(slide, Inches(0.56), Inches(1.04), Inches(4.25), Inches(0.82), title, size=21, color=BROWN, bold=True)
    tbox(slide, Inches(0.56), Inches(1.84), Inches(4.05), Inches(0.48), subtitle, size=10.3, color=MUTED)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.56), Inches(2.55), Inches(4.12), Inches(4.75))
    panel.fill.solid(); panel.fill.fore_color.rgb = WHITE; panel.line.color.rgb = LINE
    y = 2.92
    for b in bullets:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.84), Inches(y+0.04), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = GOLD; dot.line.fill.background()
        tbox(slide, Inches(1.02), Inches(y), Inches(3.2), Inches(0.35), b, size=10.4, color=CHARCOAL)
        y += 0.72
    footer(slide)
    page_badge(slide, page_no)
    return slide


def cta_slide(prs, photo_key, title, subtitle, cta_text, page_no=6):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(proc[photo_key]), 0, 0, width=prs.slide_width, height=prs.slide_height)
    over = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    over.fill.solid(); over.fill.fore_color.rgb = RGBColor(18,18,18); over.fill.transparency = 48; over.line.fill.background()
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.72), Inches(0.14), Inches(8.85))
    side.fill.solid(); side.fill.fore_color.rgb = GOLD; side.line.fill.background()
    tag(slide, Inches(0.72), Inches(0.72), Inches(2.15), 'ATENDIMENTO VITAL SLIM', fill=WHITE, tc=BROWN)
    tbox(slide, Inches(0.72), Inches(1.95), Inches(5.15), Inches(1.55), title, size=27, color=WHITE, bold=True)
    tbox(slide, Inches(0.74), Inches(3.9), Inches(5.0), Inches(0.86), subtitle, size=14.1, color=WHITE)
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.74), Inches(5.38), Inches(3.12), Inches(0.58))
    btn.fill.solid(); btn.fill.fore_color.rgb = GOLD; btn.line.fill.background()
    tf = btn.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = cta_text
    r.font.name = 'Aptos'; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
    footer(slide, dark=True)
    page_badge(slide, page_no, dark=True)
    return slide

# -------- GLYNAC --------
prs = new_prs()
cover_slide(prs, '257', 'DRA. DANIELY FREITAS', 'VITAL SLIM', 'GLYNAC\nNA ROTINA', 'Conteúdo clínico com sofisticação visual, dados objetivos e aplicação prática.')
science_slide(prs, '254', 2, 'PROVA 01', 'GLUTATIONA\nEM ALTA', 'Quando bem indicado, o protocolo tende a melhorar o ambiente redox.', [
    ('Síntese de GSH', '~120%', 'produção antioxidante mais robusta'),
    ('GSH total', '~55%', 'melhora da reserva celular'),
    ('Relação GSH/GSSG', '~125%', 'redox mais favorável'),
    ('Radicais reativos', '~50%', 'menos estresse oxidativo'),
], ref='Dados editoriais baseados no estudo de Lai et al., 2024.')
science_slide(prs, '255', 3, 'PROVA 02', 'METABOLISMO\nMAIS FLEXÍVEL', 'A leitura clínica sugere resposta metabólica mais eficiente.', [
    ('NEFA em jejum', '~50%', 'melhor mobilização de gordura'),
    ('NEFA alimentada', '~70%', 'resposta energética mais forte'),
    ('Carboidratos', '~20%', 'uso pós-refeição mais eficiente'),
    ('Supressão de NEFA', '~135%', 'sinal metabólico mais ajustado'),
], ref='Resumo visual para educação; sempre individualizar indicação.')
science_slide(prs, '256', 4, 'PROVA 03', 'COMPOSIÇÃO E\nFORÇA EM 14 DIAS', 'O objetivo é performance com critério, não modismo.', [
    ('Gordura total', '21,6 → 20,0', 'redução objetiva'),
    ('Massa magra', '60,1 → 61,0', 'preservação com leve ganho'),
    ('Força dominante', '35 → 37', 'melhora funcional'),
    ('Força não dominante', '31 → 34', 'ganho bilateral consistente'),
], ref='Resultados variam; protocolo exige avaliação médica.')
practice_slide(prs, '274', 'COMO EU LEVO\nISSO PARA A CLÍNICA', 'Nada de suplemento aleatório: tudo começa com contexto e segurança.', [
    'Avalio glutationa, inflamação e contexto metabólico',
    'Monitoro homocisteína e função hepática',
    'Ajusto dose e timing conforme resposta clínica',
    'Integro com nutrição e composição corporal',
    'Reavalio sinais, sintomas e exames ao longo do processo',
])
cta_slide(prs, '271', 'NÃO É SOBRE\nTOMAR MAIS.', 'É sobre indicar melhor, monitorar de perto e ajustar com estratégia.', 'Salve este post')
path_glynac = OUT / '2026-04-01-glynac-dra-daniely-vitalslim.pptx'
prs.save(path_glynac)
print(path_glynac)

# -------- CAPILAR --------
prs = new_prs()
cover_slide(prs, '249', 'DRA. DANIELY FREITAS', 'VITAL SLIM', 'EMAGRECER SEM\nPERDER CABELO?', 'Sim — quando o processo respeita nutrientes, hormônios e velocidade de resposta.')
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, prs)
slide.shapes.add_picture(str(proc['273']), Inches(5.18), Inches(0.78), width=Inches(2.66), height=Inches(8.05))
overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.18), Inches(0.78), Inches(2.66), Inches(8.05)); overlay.fill.solid(); overlay.fill.fore_color.rgb = RGBColor(20,20,20); overlay.fill.transparency = 52; overlay.line.fill.background()
tag(slide, Inches(0.56), Inches(0.64), Inches(1.58), 'PILAR 01', fill=BROWN, tc=WHITE)
tbox(slide, Inches(0.56), Inches(1.04), Inches(4.2), Inches(0.82), 'A QUEDA NÃO\nSURGE DO NADA', size=22, color=BROWN, bold=True)
tbox(slide, Inches(0.56), Inches(1.84), Inches(4.05), Inches(0.5), 'Na prática, o cabelo costuma sofrer quando o emagrecimento vira agressão metabólica.', size=10.2, color=MUTED)
card(slide, Inches(0.56), Inches(2.55), Inches(1.78), Inches(1.34), 'Proteína', 'Ingestão baixa reduz matéria-prima para fios mais fortes.')
card(slide, Inches(2.57), Inches(2.55), Inches(1.78), Inches(1.34), 'Ferritina', 'Reserva inadequada pode aparecer como afinamento e queda difusa.')
card(slide, Inches(0.56), Inches(4.1), Inches(1.78), Inches(1.34), 'Estresse', 'Déficit excessivo e rotina caótica aumentam a chance de eflúvio.')
card(slide, Inches(2.57), Inches(4.1), Inches(1.78), Inches(1.34), 'Velocidade', 'Emagrecer rápido demais costuma cobrar um preço biológico.')
refbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.56), Inches(5.82), Inches(3.9), Inches(0.46)); refbox.fill.solid(); refbox.fill.fore_color.rgb = BEIGE; refbox.line.fill.background()
tbox(slide, Inches(0.76), Inches(5.93), Inches(3.45), Inches(0.18), 'Referência editorial: eflúvio telógeno e perdas por estresse metabólico.', size=8.1, color=BROWN)
footer(slide); page_badge(slide, 2)
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, prs)
slide.shapes.add_picture(str(proc['257']), Inches(5.18), Inches(0.78), width=Inches(2.66), height=Inches(8.05))
overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.18), Inches(0.78), Inches(2.66), Inches(8.05)); overlay.fill.solid(); overlay.fill.fore_color.rgb = RGBColor(20,20,20); overlay.fill.transparency = 52; overlay.line.fill.background()
tag(slide, Inches(0.56), Inches(0.64), Inches(1.58), 'PILAR 02', fill=BROWN, tc=WHITE)
tbox(slide, Inches(0.56), Inches(1.04), Inches(4.2), Inches(0.82), 'NO VITAL SLIM\nOLHAMOS A CAUSA', size=22, color=BROWN, bold=True)
tbox(slide, Inches(0.56), Inches(1.84), Inches(4.05), Inches(0.5), 'O foco não é esconder a queda. É corrigir o terreno que levou a ela.', size=10.2, color=MUTED)
card(slide, Inches(0.56), Inches(2.55), Inches(1.78), Inches(1.34), 'Exames', 'Ferritina, B12, vitamina D, proteínas e contexto hormonal.')
card(slide, Inches(2.57), Inches(2.55), Inches(1.78), Inches(1.34), 'Nutrição', 'Ajustes de proteína, micronutrientes e estratégia alimentar.')
card(slide, Inches(0.56), Inches(4.1), Inches(1.78), Inches(1.34), 'Medicamentos', 'Revisão do protocolo para emagrecer sem sacrificar saúde capilar.')
card(slide, Inches(2.57), Inches(4.1), Inches(1.78), Inches(1.34), 'Ritmo', 'A velocidade do processo também é tratamento.')
refbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.56), Inches(5.82), Inches(3.9), Inches(0.46)); refbox.fill.solid(); refbox.fill.fore_color.rgb = BEIGE; refbox.line.fill.background()
tbox(slide, Inches(0.76), Inches(5.93), Inches(3.45), Inches(0.18), 'Estratégia clínica: emagrecer preservando metabolismo e fios.', size=8.1, color=BROWN)
footer(slide); page_badge(slide, 3)
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, prs)
slide.shapes.add_picture(str(proc['271']), Inches(5.18), Inches(0.78), width=Inches(2.66), height=Inches(8.05))
overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.18), Inches(0.78), Inches(2.66), Inches(8.05)); overlay.fill.solid(); overlay.fill.fore_color.rgb = RGBColor(20,20,20); overlay.fill.transparency = 52; overlay.line.fill.background()
tag(slide, Inches(0.56), Inches(0.64), Inches(1.58), 'PILAR 03', fill=BROWN, tc=WHITE)
tbox(slide, Inches(0.56), Inches(1.04), Inches(4.25), Inches(0.82), 'O OBJETIVO É\nEMAGRECER BEM', size=22, color=BROWN, bold=True)
tbox(slide, Inches(0.56), Inches(1.84), Inches(4.05), Inches(0.5), 'Resultado bonito é aquele que melhora corpo, energia e autoestima ao mesmo tempo.', size=10.2, color=MUTED)
card(slide, Inches(0.56), Inches(2.55), Inches(1.78), Inches(1.34), 'Menos perda', 'menos risco de afinamento e queda difusa.')
card(slide, Inches(2.57), Inches(2.55), Inches(1.78), Inches(1.34), 'Mais coerência', 'emagrecimento alinhado com saúde real.')
card(slide, Inches(0.56), Inches(4.1), Inches(1.78), Inches(1.34), 'Mais confiança', 'o espelho passa a acompanhar o processo.')
card(slide, Inches(2.57), Inches(4.1), Inches(1.78), Inches(1.34), 'Mais resultado', 'massa, nutrientes e cabelos mais preservados.')
refbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.56), Inches(5.82), Inches(3.9), Inches(0.46)); refbox.fill.solid(); refbox.fill.fore_color.rgb = BEIGE; refbox.line.fill.background()
tbox(slide, Inches(0.76), Inches(5.93), Inches(3.45), Inches(0.18), 'Peça editorial para educação e posicionamento médico.', size=8.1, color=BROWN)
footer(slide); page_badge(slide, 4)
practice_slide(prs, '274', 'O QUE EU FAÇO\nNA PRÁTICA', 'O protocolo certo protege o emagrecimento sem negligenciar o cabelo.', [
    'Mapeio exames e história clínica antes de acelerar o processo',
    'Corrijo proteína, ferro, vitaminas e sinais de carência',
    'Ajusto medicações e metas de velocidade',
    'Monitoro sintomas, composição corporal e resposta capilar',
    'Personalizo a estratégia conforme cada paciente',
])
cta_slide(prs, '248', 'EMAGRECER E\nPRESERVAR OS FIOS', 'É possível quando existe diagnóstico, estratégia e acompanhamento.', 'Envie CAPILAR')
path_capilar = OUT / '2026-04-01-emagrecer-sem-perder-cabelo-dra-daniely-vitalslim.pptx'
prs.save(path_capilar)
print(path_capilar)

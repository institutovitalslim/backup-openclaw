from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
from pathlib import Path

BASE = Path('/root/.openclaw/workspace')
OUT = BASE / 'deliverables'
TMP = OUT / 'tmp_glynac'
TMP.mkdir(parents=True, exist_ok=True)

photos = [
    Path('/root/.openclaw/media/inbound/file_186---e75a5ed2-8053-4cd0-bcdd-2b5ef5a70bc1.jpg'),
    Path('/root/.openclaw/media/inbound/file_187---3825409d-f26b-49ca-b763-4884eb4c4091.jpg'),
    Path('/root/.openclaw/media/inbound/file_188---379cfcdb-eb2d-4c8b-a2c8-487e1823d226.jpg'),
    Path('/root/.openclaw/media/inbound/file_189---7c6d471e-d5e7-4d9d-aa3e-a58397dc6a84.jpg'),
    Path('/root/.openclaw/media/inbound/file_190---221a1383-edf3-427e-bcef-1e6be5de8498.jpg'),
    Path('/root/.openclaw/media/inbound/file_191---76af0839-730b-44e9-bacd-b20082c8c29e.jpg'),
]

# palette
BG = RGBColor(246, 244, 239)
GREEN = RGBColor(14, 92, 76)
GREEN_DARK = RGBColor(10, 55, 45)
SAGE = RGBColor(216, 229, 223)
TEXT = RGBColor(25, 29, 27)
MUTED = RGBColor(90, 101, 96)
GOLD = RGBColor(197, 165, 102)
WHITE = RGBColor(255, 255, 255)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(210, 220, 215)

prs = Presentation()
prs.slide_width = Inches(8.333)
prs.slide_height = Inches(10.417)
SW = prs.slide_width
SH = prs.slide_height


def crop_fill(src, dst, size, dark=False, right_shadow=False):
    img = Image.open(src).convert('RGB')
    tw, th = size
    src_ratio = img.width / img.height
    dst_ratio = tw / th
    if src_ratio > dst_ratio:
        new_h = img.height
        new_w = int(new_h * dst_ratio)
        left = (img.width - new_w) // 2
        img = img.crop((left, 0, left + new_w, img.height))
    else:
        new_w = img.width
        new_h = int(new_w / dst_ratio)
        top = (img.height - new_h) // 2
        img = img.crop((0, top, img.width, top + new_h))
    img = img.resize((tw, th), Image.Resampling.LANCZOS)
    if dark or right_shadow:
        overlay = Image.new('RGBA', (tw, th), (0, 0, 0, 0))
        draw = ImageDraw.Draw(overlay)
        if dark:
            for y in range(th):
                alpha = int(140 * (y / th))
                draw.line((0, y, tw, y), fill=(7, 15, 13, alpha))
        if right_shadow:
            for x in range(tw):
                alpha = int(180 * (x / tw))
                draw.line((0, 0, x, th), fill=(5, 15, 12, alpha // 3))
            for x in range(tw):
                alpha = int(190 * ((tw - x) / tw))
                draw.line((x, 0, x, th), fill=(5, 15, 12, alpha))
        img = Image.alpha_composite(img.convert('RGBA'), overlay).convert('RGB')
    img.save(dst, quality=95)
    return dst


def add_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()


def textbox(slide, left, top, width, height, text='', size=18, color=TEXT, bold=False,
            font='Aptos', align=PP_ALIGN.LEFT, name=None, all_caps=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper() if all_caps else text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    return box


def pill(slide, left, top, width, height, text, fill=GREEN, text_color=WHITE, size=10):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    tf = shp.text_frame
    tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = 'Aptos'; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = text_color
    return shp


def page_badge(slide, num):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.55), Inches(0.45), Inches(0.38), Inches(0.38))
    c.fill.solid(); c.fill.fore_color.rgb = GREEN
    c.line.fill.background()
    tf = c.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num)
    r.font.name = 'Aptos'; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE


def footer(slide, text='@dradaniely.freitas'):
    textbox(slide, Inches(0.5), Inches(9.8), Inches(3.2), Inches(0.25), text, size=10, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(9.58), Inches(7.25), Inches(0.015))
    line.fill.solid(); line.fill.fore_color.rgb = LINE; line.line.fill.background()


def metric_card(slide, left, top, w, h, title, value, note):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = CARD
    shp.line.color.rgb = LINE
    textbox(slide, left + Inches(0.18), top + Inches(0.12), w - Inches(0.36), Inches(0.28), title, size=10, color=MUTED)
    textbox(slide, left + Inches(0.18), top + Inches(0.38), w - Inches(0.36), Inches(0.42), value, size=20, color=GREEN_DARK, bold=True)
    textbox(slide, left + Inches(0.18), top + Inches(0.82), w - Inches(0.36), Inches(0.28), note, size=9, color=TEXT)

# Prepare processed images
processed = []
for i, p in enumerate(photos, start=1):
    dst = TMP / f'crop_{i}.jpg'
    crop_fill(p, dst, (1200, 1600), dark=(i in [1, 6]))
    processed.append(dst)

# Slide 1 cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = GREEN_DARK
slide.shapes.add_picture(str(processed[0]), 0, 0, width=SW, height=SH)
shade = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
shade.fill.solid(); shade.fill.fore_color.rgb = RGBColor(8, 20, 18); shade.fill.transparency = 35
shade.line.fill.background()
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.72), Inches(0.16), Inches(8.8))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN
bar.line.fill.background()
pill(slide, Inches(0.72), Inches(0.7), Inches(2.15), Inches(0.32), 'Vital Slim • Conteúdo clínico', fill=RGBColor(255,255,255), text_color=GREEN_DARK, size=10)
textbox(slide, Inches(0.72), Inches(2.1), Inches(5.6), Inches(1.5), 'GLYNAC\nNA ROTINA', size=28, color=WHITE, bold=True)
textbox(slide, Inches(0.74), Inches(3.9), Inches(4.9), Inches(1.0), 'Glutationa, metabolismo e performance com leitura clínica sofisticada.', size=15, color=WHITE)
pill(slide, Inches(0.74), Inches(5.0), Inches(2.55), Inches(0.32), 'Dra. Daniely Freitas', fill=GREEN, text_color=WHITE, size=11)
textbox(slide, Inches(0.74), Inches(8.95), Inches(4.0), Inches(0.3), '@dradaniely.freitas', size=10, color=WHITE)
page_badge(slide, 1)

# Generic science slide creator

def science_slide(slide_num, photo_idx, eyebrow, title, subtitle, metrics, ref='Lai et al., 2024'):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    photo_w = Inches(2.85)
    slide.shapes.add_picture(str(processed[photo_idx]), Inches(5.12), Inches(0.72), width=photo_w, height=Inches(8.2))
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.12), Inches(0.72), photo_w, Inches(8.2))
    overlay.fill.solid(); overlay.fill.fore_color.rgb = RGBColor(9, 28, 24); overlay.fill.transparency = 40
    overlay.line.fill.background()
    pill(slide, Inches(0.55), Inches(0.6), Inches(2.15), Inches(0.3), eyebrow, fill=SAGE, text_color=GREEN_DARK, size=10)
    textbox(slide, Inches(0.55), Inches(1.05), Inches(4.05), Inches(0.8), title, size=22, color=GREEN_DARK, bold=True)
    textbox(slide, Inches(0.55), Inches(1.76), Inches(4.1), Inches(0.55), subtitle, size=11, color=MUTED)
    x1, x2 = Inches(0.55), Inches(2.58)
    y = Inches(2.45)
    w, h = Inches(1.82), Inches(1.2)
    for idx, (t, v, n) in enumerate(metrics):
        metric_card(slide, x1 if idx % 2 == 0 else x2, y + Inches(1.4) * (idx // 2), w, h, t, v, n)
    textbox(slide, Inches(0.55), Inches(8.9), Inches(3.8), Inches(0.25), f'Referência científica: {ref}', size=8, color=MUTED)
    footer(slide)
    page_badge(slide, slide_num)
    return slide

science_slide(2, 1, 'Prova científica 1', 'GLUTATIONA EM ALTA', 'Marcadores centrais respondem com clareza quando há estratégia e critério.', [
    ('Síntese de GSH', '~120%', 'ganho expressivo de produção'),
    ('GSH total', '~55%', 'reserva antioxidante ampliada'),
    ('Relação GSH/GSSG', '~125%', 'ambiente redox mais favorável'),
    ('Radicais reativos', '~50%', 'redução do estresse oxidativo'),
])

science_slide(3, 2, 'Prova científica 2', 'METABOLISMO MAIS FLEXÍVEL', 'Melhora da oxidação de substratos e resposta metabólica mais eficiente.', [
    ('NEFA em jejum', '~50%', 'maior mobilização lipídica'),
    ('NEFA alimentada', '~70%', 'resposta energética superior'),
    ('Carboidratos', '~20%', 'melhor uso em estado alimentado'),
    ('Supressão de NEFA', '~135%', 'sinal metabólico mais ajustado'),
])

science_slide(4, 3, 'Prova científica 3', 'COMPOSIÇÃO E FORÇA', 'Mudanças rápidas ficam mais elegantes quando o protocolo respeita o organismo.', [
    ('Gordura total', '21,6 → 20,0 kg', 'queda objetiva em 14 dias'),
    ('Massa magra', '60,1 → 61,0 kg', 'preservação com ganho leve'),
    ('Força dominante', '35 → 37', 'melhora funcional observável'),
    ('Força não dominante', '31 → 34', 'resposta global mais consistente'),
])

# Slide 5 minha prática
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide.shapes.add_picture(str(processed[4]), Inches(4.95), Inches(0.75), width=Inches(2.95), height=Inches(8.15))
overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.95), Inches(0.75), Inches(2.95), Inches(8.15))
overlay.fill.solid(); overlay.fill.fore_color.rgb = RGBColor(9, 28, 24); overlay.fill.transparency = 46
overlay.line.fill.background()
pill(slide, Inches(0.55), Inches(0.6), Inches(1.75), Inches(0.3), 'Minha prática', fill=SAGE, text_color=GREEN_DARK, size=10)
textbox(slide, Inches(0.55), Inches(1.05), Inches(3.9), Inches(0.8), 'COMO EU APLICO\nNA CLÍNICA', size=21, color=GREEN_DARK, bold=True)
textbox(slide, Inches(0.55), Inches(1.86), Inches(4.0), Inches(0.55), 'Nada de protocolo aleatório: decisão clínica, monitoramento e personalização.', size=11, color=MUTED)
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(2.55), Inches(3.95), Inches(4.95))
card.fill.solid(); card.fill.fore_color.rgb = CARD
card.line.color.rgb = LINE
bullets = [
    'Avaliação de glutationa e contexto inflamatório',
    'Homocisteína e função hepática monitoradas',
    'NAC e glicina fracionados conforme tolerância',
    'Insulina e composição corporal acompanhadas',
    'Ajustes integrados com TRH e nutrição personalizada',
]
y = Inches(2.9)
for b in bullets:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.82), y + Inches(0.08), Inches(0.1), Inches(0.1))
    dot.fill.solid(); dot.fill.fore_color.rgb = GREEN
    dot.line.fill.background()
    textbox(slide, Inches(1.0), y, Inches(3.15), Inches(0.42), b, size=11, color=TEXT)
    y += Inches(0.75)
textbox(slide, Inches(0.82), Inches(6.85), Inches(3.0), Inches(0.4), 'Estratégia clínica com sofisticação e segurança.', size=10, color=GREEN_DARK, bold=True)
footer(slide)
page_badge(slide, 5)

# Slide 6 CTA
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = GREEN_DARK
slide.shapes.add_picture(str(processed[5]), 0, 0, width=SW, height=SH)
shade = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
shade.fill.solid(); shade.fill.fore_color.rgb = RGBColor(6, 16, 14); shade.fill.transparency = 42
shade.line.fill.background()
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.78), Inches(0.16), Inches(8.6))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN
bar.line.fill.background()
pill(slide, Inches(0.74), Inches(0.72), Inches(2.4), Inches(0.32), 'Vital Slim • Lauro de Freitas', fill=RGBColor(255,255,255), text_color=GREEN_DARK, size=10)
textbox(slide, Inches(0.74), Inches(2.0), Inches(5.2), Inches(1.3), 'NÃO É\nSUPLEMENTO ALEATÓRIO.', size=27, color=WHITE, bold=True)
textbox(slide, Inches(0.74), Inches(4.15), Inches(4.8), Inches(0.9), 'É protocolo com leitura clínica, ajuste fino e acompanhamento real.', size=15, color=WHITE)
cta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.74), Inches(5.45), Inches(3.2), Inches(0.6))
cta.fill.solid(); cta.fill.fore_color.rgb = GREEN
cta.line.fill.background()
tf = cta.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'Agende sua avaliação integral'; r.font.name='Aptos'; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=WHITE
textbox(slide, Inches(0.76), Inches(8.9), Inches(4.2), Inches(0.3), '@dradaniely.freitas', size=10, color=WHITE)
page_badge(slide, 6)

out_path = OUT / 'glynac-dra-daniely-vitalslim.pptx'
prs.save(out_path)
print(out_path)

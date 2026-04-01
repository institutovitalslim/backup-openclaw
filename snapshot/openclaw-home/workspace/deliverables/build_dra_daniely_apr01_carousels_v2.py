from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
from pathlib import Path

OUT = Path('/root/.openclaw/workspace/deliverables')
TMP = OUT / 'tmp_apr01_carousels_v2'
TMP.mkdir(parents=True, exist_ok=True)

GOLD = RGBColor(159, 136, 68)
BROWN = RGBColor(87, 71, 23)
CREAM = RGBColor(247, 243, 236)
CHAR = RGBColor(24, 24, 24)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(118, 104, 82)
LINE = RGBColor(226, 218, 206)
SOFT = RGBColor(236, 230, 219)

photos = {
    '248': Path('/root/.openclaw/media/inbound/file_248---754c4855-061a-4e03-ae26-08b4a16b1da4.jpg'),
    '249': Path('/root/.openclaw/media/inbound/file_249---17318fde-dda9-4335-80a3-c22b8a774885.jpg'),
    '254': Path('/root/.openclaw/media/inbound/file_254---dfb15256-8d77-41ea-8aab-79b84f76357b.jpg'),
    '255': Path('/root/.openclaw/media/inbound/file_255---28c713b0-b988-4762-a1fe-0b96118a0f87.jpg'),
    '256': Path('/root/.openclaw/media/inbound/file_256---1133538b-77c4-4ae8-99a8-d8dedf3ccfb6.jpg'),
    '257': Path('/root/.openclaw/media/inbound/file_257---f0a37b69-9ccd-44c0-a981-4483bd299d9d.jpg'),
    '271': Path('/root/.openclaw/media/inbound/file_271---3696783f-b5d9-46dc-ba4d-3aedabbbbf0c.jpg'),
    '273': Path('/root/.openclaw/media/inbound/file_273---3f0ce689-e3db-48b7-85ca-971215f608f1.jpg'),
    '274': Path('/root/.openclaw/media/inbound/file_274---87a94e30-217c-4ff8-bb04-c012647be574.jpg'),
}


def crop_fill(src, dst, size):
    img = Image.open(src).convert('RGB')
    tw, th = size
    src_ratio = img.width / img.height
    dst_ratio = tw / th
    if src_ratio > dst_ratio:
        nw = int(img.height * dst_ratio)
        left = (img.width - nw) // 2
        img = img.crop((left, 0, left + nw, img.height))
    else:
        nh = int(img.width / dst_ratio)
        top = (img.height - nh) // 2
        img = img.crop((0, top, img.width, top + nh))
    img = img.resize((tw, th), Image.Resampling.LANCZOS)
    img.save(dst, quality=95)
    return dst

proc = {}
for k, p in photos.items():
    dst = TMP / f'{k}.jpg'
    crop_fill(p, dst, (1200, 1600))
    proc[k] = dst


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(8.333)
    prs.slide_height = Inches(10.417)
    return prs


def bg(slide, prs):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = CREAM
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), prs.slide_height)
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()


def text(slide, l, t, w, h, txt, size=18, color=CHAR, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.name = 'Aptos'; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box


def pill(slide, l, t, w, txt, fill=BROWN, tc=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Inches(0.3))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background()
    tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.name = 'Aptos'; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = tc


def footer(slide, dark=False):
    col = WHITE if dark else MUTED
    linec = RGBColor(110,110,110) if dark else LINE
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(9.58), Inches(7.12), Inches(0.015))
    ln.fill.solid(); ln.fill.fore_color.rgb = linec; ln.line.fill.background()
    text(slide, Inches(0.56), Inches(9.72), Inches(2.6), Inches(0.18), '@dradaniely.freitas', size=9, color=col)
    arr = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(7.12), Inches(9.67), Inches(0.17), Inches(0.11))
    arr.fill.solid(); arr.fill.fore_color.rgb = col; arr.line.fill.background()


def page(slide, n, dark=False):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.56), Inches(0.44), Inches(0.38), Inches(0.38))
    c.fill.solid(); c.fill.fore_color.rgb = WHITE if dark else GOLD; c.line.fill.background()
    tf = c.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.name='Aptos'; r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb = BROWN if dark else WHITE


def photo_panel(slide, path, l, t, w, h, dark_overlay=False):
    slide.shapes.add_picture(str(path), l, t, width=w, height=h)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    panel.fill.background(); panel.line.color.rgb = GOLD; panel.line.width = Pt(1)
    if dark_overlay:
        ov = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        ov.fill.solid(); ov.fill.fore_color.rgb = RGBColor(20,20,20); ov.fill.transparency = 55; ov.line.fill.background()


def card_metric(slide, l, t, w, h, title, value, note):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = WHITE; s.line.color.rgb = LINE
    text(slide, l+Inches(0.14), t+Inches(0.10), w-Inches(0.28), Inches(0.2), title.upper(), size=8.5, color=MUTED)
    text(slide, l+Inches(0.14), t+Inches(0.30), w-Inches(0.28), Inches(0.3), value, size=15, color=BROWN, bold=True)
    text(slide, l+Inches(0.14), t+Inches(0.60), w-Inches(0.28), Inches(0.28), note, size=8.2, color=CHAR)


def cover(prs, photo_key, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, prs)
    # photo takes real space on right
    photo_panel(slide, proc[photo_key], Inches(4.8), Inches(0.6), Inches(2.85), Inches(8.95), dark_overlay=False)
    tint = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(0.6), Inches(2.85), Inches(8.95))
    tint.fill.solid(); tint.fill.fore_color.rgb = RGBColor(0,0,0); tint.fill.transparency = 72; tint.line.fill.background()
    pill(slide, Inches(0.56), Inches(0.62), Inches(2.28), 'DRA. DANIELY FREITAS', fill=WHITE, tc=BROWN)
    pill(slide, Inches(2.92), Inches(0.62), Inches(1.45), 'VITAL SLIM', fill=GOLD, tc=WHITE)
    text(slide, Inches(0.56), Inches(2.0), Inches(3.8), Inches(1.6), title, size=25, color=BROWN, bold=True)
    text(slide, Inches(0.58), Inches(4.12), Inches(3.65), Inches(1.05), subtitle, size=13.2, color=CHAR)
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.56), Inches(6.0), Inches(3.35), Inches(0.55))
    quote.fill.solid(); quote.fill.fore_color.rgb = SOFT; quote.line.fill.background()
    text(slide, Inches(0.74), Inches(6.17), Inches(2.95), Inches(0.18), 'Ciência, sofisticação e autoridade clínica.', size=9.5, color=BROWN, bold=True)
    footer(slide)
    page(slide, 1)


def evidence(prs, photo_key, page_no, chip, title, subtitle, metrics, ref):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, prs)
    photo_panel(slide, proc[photo_key], Inches(5.0), Inches(1.0), Inches(2.65), Inches(7.75), dark_overlay=False)
    pill(slide, Inches(0.56), Inches(0.62), Inches(1.45), chip, fill=BROWN, tc=WHITE)
    text(slide, Inches(0.56), Inches(1.0), Inches(4.1), Inches(0.82), title, size=20.5, color=BROWN, bold=True)
    text(slide, Inches(0.56), Inches(1.78), Inches(4.0), Inches(0.45), subtitle, size=9.8, color=MUTED)
    positions=[(0.56,2.5),(2.56,2.5),(0.56,3.95),(2.56,3.95)]
    for (x,y),(a,b,c) in zip(positions, metrics):
        card_metric(slide, Inches(x), Inches(y), Inches(1.78), Inches(1.08), a, b, c)
    refb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.56), Inches(5.55), Inches(3.95), Inches(0.45))
    refb.fill.solid(); refb.fill.fore_color.rgb = SOFT; refb.line.fill.background()
    text(slide, Inches(0.76), Inches(5.68), Inches(3.45), Inches(0.18), ref, size=8, color=BROWN)
    footer(slide)
    page(slide, page_no)


def practice(prs, photo_key, title, subtitle, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, prs)
    photo_panel(slide, proc[photo_key], Inches(4.95), Inches(0.88), Inches(2.7), Inches(8.0), dark_overlay=False)
    pill(slide, Inches(0.56), Inches(0.62), Inches(1.8), 'MINHA PRÁTICA', fill=GOLD, tc=WHITE)
    text(slide, Inches(0.56), Inches(1.0), Inches(4.0), Inches(0.82), title, size=20.5, color=BROWN, bold=True)
    text(slide, Inches(0.56), Inches(1.78), Inches(4.0), Inches(0.45), subtitle, size=10, color=MUTED)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.56), Inches(2.45), Inches(4.1), Inches(4.95))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = LINE
    y=2.82
    for b in bullets:
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.82), Inches(y+0.03), Inches(0.1), Inches(0.1))
        d.fill.solid(); d.fill.fore_color.rgb = GOLD; d.line.fill.background()
        text(slide, Inches(1.02), Inches(y), Inches(3.3), Inches(0.34), b, size=10.1, color=CHAR)
        y += 0.72
    footer(slide)
    page(slide, 5)


def cta(prs, photo_key, title, subtitle, cta_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = CHAR
    photo_panel(slide, proc[photo_key], Inches(4.78), Inches(0.62), Inches(2.88), Inches(8.95), dark_overlay=False)
    dark = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.78), Inches(0.62), Inches(2.88), Inches(8.95))
    dark.fill.solid(); dark.fill.fore_color.rgb = RGBColor(0,0,0); dark.fill.transparency = 58; dark.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.72), Inches(0.14), Inches(8.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()
    pill(slide, Inches(0.72), Inches(0.72), Inches(2.1), 'ATENDIMENTO VITAL SLIM', fill=WHITE, tc=BROWN)
    text(slide, Inches(0.72), Inches(2.0), Inches(3.7), Inches(1.35), title, size=23.5, color=WHITE, bold=True)
    text(slide, Inches(0.74), Inches(3.8), Inches(3.7), Inches(0.95), subtitle, size=12.8, color=WHITE)
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.74), Inches(5.5), Inches(2.7), Inches(0.58))
    btn.fill.solid(); btn.fill.fore_color.rgb = GOLD; btn.line.fill.background()
    tf = btn.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = cta_text; r.font.name='Aptos'; r.font.size=Pt(14); r.font.bold=True; r.font.color.rgb=WHITE
    footer(slide, dark=True)
    page(slide, 6, dark=True)

# GLYNAC corrected
prs = new_prs()
cover(prs, '257', 'GLYNAC\nNA ROTINA', 'Dados e aplicação clínica em um layout com mais autoridade visual e presença da Dra. Daniely.')
evidence(prs, '254', 2, 'PROVA 01', 'GLUTATIONA\nEM ALTA', 'Quando bem indicado, o protocolo pode melhorar o ambiente redox.', [
    ('Síntese de GSH', '~120%', 'produção antioxidante mais forte'),
    ('GSH total', '~55%', 'reserva celular ampliada'),
    ('GSH/GSSG', '~125%', 'redox mais favorável'),
    ('Radicais', '~50%', 'menos estresse oxidativo'),
], 'Base editorial: Lai et al., 2024.')
evidence(prs, '255', 3, 'PROVA 02', 'METABOLISMO\nMAIS FLEXÍVEL', 'A leitura clínica sugere resposta metabólica mais eficiente.', [
    ('NEFA jejum', '~50%', 'melhor mobilização lipídica'),
    ('NEFA alimentada', '~70%', 'resposta energética maior'),
    ('Carboidratos', '~20%', 'uso pós-refeição mais eficaz'),
    ('Supressão NEFA', '~135%', 'sinal metabólico ajustado'),
], 'Resumo visual para educação; individualizar sempre.')
evidence(prs, '256', 4, 'PROVA 03', 'COMPOSIÇÃO E\nFORÇA EM 14 DIAS', 'O objetivo é performance com critério — não modismo.', [
    ('Gordura total', '21,6 → 20,0', 'queda objetiva'),
    ('Massa magra', '60,1 → 61,0', 'preservação com leve ganho'),
    ('Força dom.', '35 → 37', 'melhora funcional'),
    ('Força não dom.', '31 → 34', 'ganho consistente'),
], 'Resultados variam; avaliação médica é indispensável.')
practice(prs, '274', 'COMO ISSO ENTRA\nNA MINHA PRÁTICA', 'Nada de suplemento aleatório. Tudo parte de contexto, exames e segurança.', [
    'Avalio glutationa, inflamação e contexto metabólico',
    'Monitoro homocisteína e função hepática',
    'Ajusto dose e timing conforme resposta clínica',
    'Integro com nutrição e composição corporal',
    'Reavalio sinais, sintomas e exames ao longo do processo',
])
cta(prs, '271', 'NÃO É SOBRE\nTOMAR MAIS.', 'É sobre indicar melhor, acompanhar de perto e ajustar com estratégia.', 'Salve este post')
path1 = OUT / '2026-04-01-glynac-dra-daniely-vitalslim-v2.pptx'
prs.save(path1)
print(path1)

# Hair corrected
prs = new_prs()
cover(prs, '271', 'EMAGRECER SEM\nPERDER CABELO?', 'Sim — quando o processo respeita nutrientes, hormônios e velocidade de resposta.')
evidence(prs, '273', 2, 'PILAR 01', 'A QUEDA NÃO\nSURGE DO NADA', 'O cabelo costuma sofrer quando o emagrecimento vira agressão metabólica.', [
    ('Proteína', 'base', 'menos matéria-prima para fios fortes'),
    ('Ferritina', 'alerta', 'queda difusa pede investigação'),
    ('Estresse', 'alto', 'mais risco de eflúvio telógeno'),
    ('Velocidade', 'excesso', 'resultado rápido também cobra preço'),
], 'Base editorial: eflúvio telógeno e estresse metabólico.')
evidence(prs, '257', 3, 'PILAR 02', 'NO VITAL SLIM\nOLHAMOS A CAUSA', 'O foco não é esconder a queda. É corrigir o terreno que levou a ela.', [
    ('Exames', 'mapa', 'ferritina, B12, vitamina D, proteínas'),
    ('Nutrição', 'ajuste', 'proteína e micronutrientes certos'),
    ('Medicamentos', 'revisão', 'emagrecer sem sacrificar saúde capilar'),
    ('Ritmo', 'critério', 'a velocidade também é tratamento'),
], 'Peça educativa com abordagem clínica integrada.')
evidence(prs, '249', 4, 'PILAR 03', 'O OBJETIVO É\nEMAGRECER BEM', 'Resultado bonito é aquele que melhora corpo, energia e autoestima ao mesmo tempo.', [
    ('Menos perda', 'fios', 'reduz risco de afinamento difuso'),
    ('Mais coerência', 'metabólica', 'emagrecimento alinhado à saúde'),
    ('Mais confiança', 'espelho', 'cabelo acompanha o processo'),
    ('Mais resultado', 'global', 'massa, nutrientes e fios preservados'),
], 'Educação médica com linguagem acessível.')
practice(prs, '274', 'O QUE EU FAÇO\nNA PRÁTICA', 'A estratégia certa protege o emagrecimento sem negligenciar o cabelo.', [
    'Mapeio exames e história clínica antes de acelerar o processo',
    'Corrijo proteína, ferro, vitaminas e sinais de carência',
    'Ajusto medicações e metas de velocidade',
    'Monitoro sintomas, composição corporal e resposta capilar',
    'Personalizo tudo conforme o paciente',
])
cta(prs, '248', 'EMAGRECER E\nPRESERVAR OS FIOS', 'É possível quando existe diagnóstico, estratégia e acompanhamento.', 'Envie CAPILAR')
path2 = OUT / '2026-04-01-emagrecer-sem-perder-cabelo-dra-daniely-vitalslim-v2.pptx'
prs.save(path2)
print(path2)

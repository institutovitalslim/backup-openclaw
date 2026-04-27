#!/usr/bin/env python3
"""
Document Presentation Creator
Cria apresentações .pptx com slides profissionais.
"""

import argparse
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def create_presentation(output_path, title, subtitle, slides_content, image_path=None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide de título
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title or "Apresentação"
    
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    
    # Adicionar imagem no slide de título se fornecida
    if image_path:
        try:
            slide.shapes.add_picture(image_path, Inches(11), Inches(0.5), width=Inches(2))
        except:
            pass
    
    # Slides de conteúdo
    if slides_content:
        slides = slides_content.split('---')
        for slide_text in slides:
            slide_text = slide_text.strip()
            if not slide_text:
                continue
            
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            
            lines = slide_text.split('\n')
            if lines:
                # Primeira linha como título do slide
                slide.shapes.title.text = lines[0][:50]
                
                # Resto como conteúdo
                if len(lines) > 1:
                    body_shape = slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.text = lines[1]
                    
                    for line in lines[2:]:
                        p = tf.add_paragraph()
                        p.text = line
                        p.level = 0
    
    # Slide final
    end_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(end_slide_layout)
    
    # Adicionar texto de agradecimento
    left = Inches(1)
    top = Inches(3)
    width = Inches(11.333)
    height = Inches(1.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Obrigado!"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xC9, 0xA9, 0x6E)
    
    # Salvar
    prs.save(output_path)
    print(f"✅ Apresentação criada: {output_path} ({len(prs.slides)} slides)")
    return output_path

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Criar apresentação PowerPoint')
    parser.add_argument('--output', '-o', required=True, help='Caminho de saída .pptx')
    parser.add_argument('--title', '-t', help='Título da apresentação')
    parser.add_argument('--subtitle', '-s', help='Subtítulo')
    parser.add_argument('--slides', help='Conteúdo dos slides (separados por ---)')
    parser.add_argument('--image', '-i', help='Imagem para slide de título')
    
    args = parser.parse_args()
    
    try:
        create_presentation(args.output, args.title, args.subtitle, args.slides, args.image)
    except Exception as e:
        print(f"❌ Erro: {e}", file=sys.stderr)
        sys.exit(1)
